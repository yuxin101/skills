#!/usr/bin/env python3
"""
Step 1.5: Convert PPTX slides to high-resolution PNG images using optimized PDF-first approach.

Usage:
    python scripts/pptx_to_images.py

Input:  slides/presentation.pptx
Output: output/images/slide_01.png, slide_02.png, ...

Optimized approach: PPTX → PDF → PNG for best quality
"""

import argparse
import subprocess
import sys
import tempfile
import shutil
from pathlib import Path

from utils import PROJECT_ROOT, SLIDES_DIR, IMAGES_DIR, PPTX_PATH, load_config


# ---- Load config (with fallback defaults) ----
try:
    _config = load_config()
    _image_cfg = _config.get("image", {})
except FileNotFoundError:
    _image_cfg = {}

# Quality presets (from config.json → image section)
PRESET_NORMAL = {
    "dpi": _image_cfg.get("dpi", 300),
    "width": _image_cfg.get("width", 1920),
    "height": _image_cfg.get("height", 1080),
}
PRESET_FAST = {
    "dpi": _image_cfg.get("dpi_fast", 150),
    "width": _image_cfg.get("width_fast", 1280),
    "height": _image_cfg.get("height_fast", 720),
}

# Output directory
OUTPUT_DIR = IMAGES_DIR


def check_libreoffice() -> bool:
    """Check if LibreOffice is available."""
    try:
        result = subprocess.run(["soffice", "--version"], capture_output=True, text=True)
        return result.returncode == 0
    except FileNotFoundError:
        return False


def pptx_to_pdf_via_libreoffice() -> Path:
    """Convert PPTX to PDF using LibreOffice headless."""
    if not check_libreoffice():
        print("Error: LibreOffice not found. Install with: brew install libreoffice")
        print("  or use: python scripts/pdf_to_images.py (requires existing PDF)")
        sys.exit(1)
    
    # Create temporary directory for conversion
    with tempfile.TemporaryDirectory() as tmpdir:
        pdf_path = Path(tmpdir) / "presentation.pdf"
        
        # Convert PPTX to PDF using LibreOffice headless
        cmd = [
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(tmpdir),
            str(PPTX_PATH)
        ]
        
        print(f"Converting PPTX to PDF via LibreOffice...")
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode != 0:
            print(f"LibreOffice conversion failed: {result.stderr}")
            sys.exit(1)
        
        if not pdf_path.exists():
            print("Error: PDF output not created by LibreOffice")
            sys.exit(1)
        
        return pdf_path


def pdf_to_images_via_pdf2image(pdf_path: Path, fast: bool = False):
    """Convert PDF to images using pdf2image (higher quality than pdftoppm)."""
    try:
        from pdf2image import convert_from_path
        from PIL import Image
    except ImportError:
        print("Error: pdf2image required for high-quality PDF conversion")
        print("Install with: pip install pdf2image")
        sys.exit(1)
    
    preset = PRESET_FAST if fast else PRESET_NORMAL
    dpi = preset["dpi"]
    target_w = preset["width"]
    target_h = preset["height"]
    
    # Ensure output directory exists
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    mode_label = "⚡ FAST" if fast else "HQ"
    print(f"Converting: {pdf_path}")
    print(f"Output dir: {OUTPUT_DIR}")
    print(f"Mode: {mode_label} | DPI: {dpi}, Target: {target_w}x{target_h}")
    print("-" * 50)
    
    # Convert PDF to images using pdf2image
    images = convert_from_path(
        str(pdf_path),
        dpi=dpi,
        fmt="png",
    )
    
    print(f"Total pages: {len(images)}")
    
    for i, img in enumerate(images, start=1):
        # Resize to target resolution with high-quality resampling
        img_resized = img.resize(
            (target_w, target_h),
            Image.LANCZOS,
        )
        
        # Save as PNG
        output_path = OUTPUT_DIR / f"slide_{i:02d}.png"
        img_resized.save(str(output_path), "PNG", optimize=True)
        
        # File size info
        size_kb = output_path.stat().st_size / 1024
        print(f"  ✅ Page {i:2d} -> {output_path.name} ({size_kb:.0f} KB)")
    
    print("-" * 50)
    print(f"Done! {len(images)} images saved to {OUTPUT_DIR}/")


def pptx_to_images_optimized(fast: bool = False):
    """Optimized PPTX to images conversion using PDF-first approach."""
    # Validate input
    if not PPTX_PATH.exists():
        print(f"Error: PPTX not found at {PPTX_PATH}")
        sys.exit(1)
    
    print("🎯 Using optimized PDF-first approach for best quality")
    print("📊 Conversion path: PPTX → PDF → PNG")
    
    # Step 1: Convert PPTX to PDF
    pdf_path = pptx_to_pdf_via_libreoffice()
    
    # Step 2: Convert PDF to images using high-quality pdf2image
    pdf_to_images_via_pdf2image(pdf_path, fast=fast)


def pptx_to_images_improved_fallback(fast: bool = False):
    """Improved fallback method with better text layout and reduced overlap."""
    try:
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("Error: python-pptx and Pillow required for fallback method")
        print("Install with: pip install python-pptx Pillow")
        sys.exit(1)
    
    preset = PRESET_FAST if fast else PRESET_NORMAL
    target_w = preset["width"]
    target_h = preset["height"]
    
    # Ensure output directory exists
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    print(f"⚠️  Using improved fallback method (LibreOffice not available)")
    print(f"📐 Target resolution: {target_w}x{target_h}")
    print("💡 Tip: Install LibreOffice for pixel-perfect rendering: brew install libreoffice")
    
    try:
        prs = Presentation(str(PPTX_PATH))
    except Exception as e:
        print(f"Error loading PPTX: {e}")
        sys.exit(1)
    
    for i, slide in enumerate(prs.slides, start=1):
        # Create slide image with better layout
        img = Image.new('RGB', (target_w, target_h), color='#ffffff')
        draw = ImageDraw.Draw(img)
        
        # Improved font handling
        try:
            title_font = ImageFont.truetype("/System/Library/Fonts/PingFang.ttc", 48)
            heading_font = ImageFont.truetype("/System/Library/Fonts/PingFang.ttc", 36)
            body_font = ImageFont.truetype("/System/Library/Fonts/PingFang.ttc", 28)
        except:
            title_font = ImageFont.load_default()
            heading_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
        
        # Improved slide layout with better spacing
        header_height = 100
        content_margin = 80
        
        # Header background
        draw.rectangle([0, 0, target_w, header_height], fill='#2c3e50')
        
        # Slide title
        slide_title = f"Slide {i}"
        title_width = draw.textlength(slide_title, font=title_font)
        title_x = (target_w - title_width) // 2
        draw.text((title_x, 30), slide_title, fill='#ffffff', font=title_font)
        
        # Extract and organize content with better spacing
        content_items = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Improved content classification
                if len(text) < 40:
                    content_items.append(('heading', text))
                elif len(text) < 200:
                    content_items.append(('bullet', text))
                else:
                    content_items.append(('paragraph', text))
        
        # Render content with improved layout
        current_y = header_height + 60
        
        for content_type, text in content_items[:8]:  # Limit to prevent overflow
            if current_y > target_h - 100:
                break
                
            if content_type == 'heading':
                # Heading with better spacing - handle multiline text
                # Use only first line for width calculation to avoid ValueError
                first_line = text.split('\n')[0] if '\n' in text else text
                text_width = draw.textlength(first_line, font=heading_font)
                text_x = content_margin
                draw.text((text_x, current_y), text, fill='#2c3e50', font=heading_font)
                current_y += 70
                
            elif content_type == 'bullet':
                # Bullet point with improved spacing
                bullet_text = f"• {text}"
                draw.text((content_margin, current_y), bullet_text, fill='#495057', font=body_font)
                current_y += 50
                
            elif content_type == 'paragraph':
                # Paragraph with text wrapping and better spacing
                max_width = target_w - (content_margin * 2)
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    test_line = ' '.join(current_line + [word])
                    # Ensure test_line is single line for textlength measurement
                    single_line = test_line.replace('\n', ' ').replace('\r', ' ')
                    if draw.textlength(single_line, font=body_font) < max_width:
                        current_line.append(word)
                    else:
                        lines.append(' '.join(current_line))
                        current_line = [word]
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                for line in lines[:4]:  # Limit lines
                    if current_y > target_h - 100:
                        break
                    draw.text((content_margin, current_y), line, fill='#495057', font=body_font)
                    current_y += 40
                
                current_y += 20
        
        # Footer
        footer_text = f"Slide {i} of {len(prs.slides)}"
        footer_width = draw.textlength(footer_text, font=body_font)
        footer_x = (target_w - footer_width) // 2
        draw.text((footer_x, target_h - 60), footer_text, fill='#6c757d', font=body_font)
        
        # Save image
        output_path = OUTPUT_DIR / f"slide_{i:02d}.png"
        img.save(str(output_path), "PNG", optimize=True, quality=95)
        
        size_kb = output_path.stat().st_size / 1024
        content_count = len(content_items)
        print(f"  🎨 Slide {i:2d} -> {output_path.name} ({size_kb:.0f} KB) [{content_count} items]")
    
    print(f"Created {len(prs.slides)} improved slide images")
    print("💡 For best quality, install LibreOffice: brew install libreoffice")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert PPTX slides to images")
    parser.add_argument("--fast", action="store_true", help="Fast mode: lower resolution (720p)")
    parser.add_argument("--fallback", action="store_true", help="Use improved fallback method")
    args = parser.parse_args()
    
    if args.fallback:
        pptx_to_images_improved_fallback(fast=args.fast)
    else:
        pptx_to_images_optimized(fast=args.fast)