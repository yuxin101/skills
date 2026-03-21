#!/usr/bin/env python3
# Part of doc2slides skill.

#!/usr/bin/env python3
"""
Extract design style from existing PowerPoint files.
Extracts colors, fonts, layout patterns to apply to new slides.

Usage:
  python extract_style.py --input presentation.pptx --output style.json
"""

import sys
import json
import argparse
from pathlib import Path
from collections import Counter

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Error: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)


def rgb_to_hex(rgb_color):
    """Convert RGBColor to hex string."""
    return f"{rgb_color.red:02x}{rgb_color.green:02x}{rgb_color.blue:02x}"


def extract_colors(prs) -> dict:
    """
    Extract color palette from presentation.
    Returns: primary, secondary, accent colors
    """
    colors = {
        'primary': None,
        'secondary': None,
        'accents': [],
        'text_colors': [],
        'backgrounds': []
    }
    
    text_color_counter = Counter()
    bg_color_counter = Counter()
    
    for slide in prs.slides[:5]:  # Analyze first 5 slides
        for shape in slide.shapes:
            # Extract text colors
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.color.rgb:
                            color_hex = rgb_to_hex(run.font.color.rgb)
                            text_color_counter[color_hex] += 1
            
            # Extract fill colors (backgrounds)
            if shape.fill.type is not None:
                try:
                    if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                        color_hex = rgb_to_hex(shape.fill.fore_color.rgb)
                        bg_color_counter[color_hex] += 1
                except:
                    pass
    
    # Get most common colors
    if text_color_counter:
        colors['text_colors'] = [c for c, _ in text_color_counter.most_common(3)]
    
    if bg_color_counter:
        colors['backgrounds'] = [c for c, _ in bg_color_counter.most_common(3)]
    
    # Extract theme colors
    try:
        theme_colors = prs.slide_master.theme.color_scheme
        if theme_colors:
            # Map theme colors
            theme_map = {
                'dk1': 'primary',
                'dk2': 'secondary',
                'accent1': 'accent1',
                'accent2': 'accent2',
                'accent3': 'accent3',
                'accent4': 'accent4'
            }
            
            for attr_name, key in theme_map.items():
                try:
                    color = getattr(theme_colors, attr_name)
                    if color and hasattr(color, 'rgb'):
                        colors[key] = rgb_to_hex(color.rgb)
                except:
                    pass
    except:
        pass
    
    # Fallback to McKinsey colors if extraction failed
    if not colors['primary']:
        colors['primary'] = '0f172a'  # Navy
    if not colors['secondary']:
        colors['secondary'] = '1e293b'
    if not colors['accents']:
        colors['accents'] = ['dc2626', 'ea580c', 'd97706', '059669']  # Red, Orange, Amber, Green
    
    return colors


def extract_fonts(prs) -> dict:
    """
    Extract font usage from presentation.
    Returns: title_font, body_font, font_sizes
    """
    fonts = {
        'title_font': None,
        'body_font': None,
        'title_size': None,
        'body_size': None
    }
    
    font_counter = Counter()
    size_counter = Counter()
    
    for slide in prs.slides[:5]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # Check font size to distinguish title/body
                    if para.font.size:
                        size_pt = para.font.size.pt
                        size_counter[size_pt] += 1
                        
                        if size_pt >= 24:  # Title
                            if para.font.name:
                                fonts['title_font'] = para.font.name
                            fonts['title_size'] = size_pt
                        elif size_pt >= 14:  # Body
                            if para.font.name:
                                fonts['body_font'] = para.font.name
                            fonts['body_size'] = size_pt
    
    # Fallback to default fonts
    if not fonts['title_font']:
        fonts['title_font'] = 'Microsoft YaHei'
    if not fonts['body_font']:
        fonts['body_font'] = 'Microsoft YaHei'
    if not fonts['title_size']:
        fonts['title_size'] = 32
    if not fonts['body_size']:
        fonts['body_size'] = 16
    
    return fonts


def extract_layout_patterns(prs) -> dict:
    """
    Extract layout patterns from presentation.
    Returns: title_position, content_area, margin
    """
    layout = {
        'title_position': {'x': 0.5, 'y': 0.3, 'width': 9.0},  # inches
        'content_area': {'x': 0.5, 'y': 1.5, 'width': 9.0, 'height': 4.5},
        'margin': 0.5,
        'slide_size': {'width': prs.slide_width.inches, 'height': prs.slide_height.inches}
    }
    
    # Analyze first few slides for layout patterns
    title_positions = []
    content_positions = []
    
    for slide in prs.slides[:3]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                x = shape.left / 914400  # EMU to inches
                y = shape.top / 914400
                w = shape.width / 914400
                h = shape.height / 914400
                
                # Assume first large text is title
                if len(text) < 100 and shape.text_frame.paragraphs[0].font.size and \
                   shape.text_frame.paragraphs[0].font.size.pt >= 24:
                    title_positions.append({'x': x, 'y': y, 'width': w})
                else:
                    content_positions.append({'x': x, 'y': y, 'width': w, 'height': h})
    
    # Average the positions
    if title_positions:
        layout['title_position'] = {
            'x': sum(p['x'] for p in title_positions) / len(title_positions),
            'y': sum(p['y'] for p in title_positions) / len(title_positions),
            'width': sum(p['width'] for p in title_positions) / len(title_positions)
        }
    
    if content_positions:
        layout['content_area'] = {
            'x': sum(p['x'] for p in content_positions) / len(content_positions),
            'y': sum(p['y'] for p in content_positions) / len(content_positions),
            'width': sum(p['width'] for p in content_positions) / len(content_positions),
            'height': sum(p['height'] for p in content_positions) / len(content_positions)
        }
    
    return layout


def extract_style(pptx_path: str) -> dict:
    """
    Extract complete style from PowerPoint file.
    """
    if not HAS_PPTX:
        return get_default_style()
    
    try:
        prs = Presentation(pptx_path)
        
        style = {
            'name': Path(pptx_path).stem,
            'colors': extract_colors(prs),
            'fonts': extract_fonts(prs),
            'layout': extract_layout_patterns(prs),
            'source_file': str(pptx_path)
        }
        
        return style
    except Exception as e:
        print(f"Warning: Failed to extract style from {pptx_path}: {e}", file=sys.stderr)
        return get_default_style()


def get_default_style() -> dict:
    """
    Return default McKinsey/BCG style.
    """
    return {
        'name': 'default',
        'colors': {
            'primary': '0f172a',
            'secondary': '1e293b',
            'accents': ['dc2626', 'ea580c', 'd97706', '059669'],
            'text_colors': ['0f172a', '475569', '94a3b8'],
            'backgrounds': ['ffffff', 'f1f5f9']
        },
        'fonts': {
            'title_font': 'Microsoft YaHei',
            'body_font': 'Microsoft YaHei',
            'title_size': 32,
            'body_size': 16
        },
        'layout': {
            'title_position': {'x': 0.5, 'y': 0.3, 'width': 9.0},
            'content_area': {'x': 0.5, 'y': 1.5, 'width': 9.0, 'height': 4.5},
            'margin': 0.5,
            'slide_size': {'width': 13.33, 'height': 7.5}
        }
    }


def main():
    parser = argparse.ArgumentParser(description='Extract style from PowerPoint')
    parser.add_argument('--input', required=True, help='Input PPTX file')
    parser.add_argument('--output', help='Output JSON file')
    
    args = parser.parse_args()
    
    style = extract_style(args.input)
    
    output_json = json.dumps(style, ensure_ascii=False, indent=2)
    
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(output_json)
        print(f"✓ Style extracted to {args.output}")
    else:
        print(output_json)


if __name__ == "__main__":
    main()
