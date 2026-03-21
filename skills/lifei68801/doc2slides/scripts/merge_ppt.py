#!/usr/bin/env python3
# Part of doc2slides skill.

#!/usr/bin/env python3
"""合并两个 PPT 文件"""

from pptx import Presentation
from pptx.util import Inches
import sys
import os

def merge_ppts(base_ppt_path, add_ppt_path, output_path):
    """将 add_ppt 的所有幻灯片追加到 base_ppt 后面"""
    
    print(f"📄 基础 PPT: {base_ppt_path}")
    print(f"📄 追加 PPT: {add_ppt_path}")
    
    # 加载两个 PPT
    base_prs = Presentation(base_ppt_path)
    add_prs = Presentation(add_ppt_path)
    
    print(f"   基础 PPT 页数: {len(base_prs.slides)}")
    print(f"   追加 PPT 页数: {len(add_prs.slides)}")
    
    # 获取基础 PPT 的幻灯片尺寸
    base_width = base_prs.slide_width
    base_height = base_prs.slide_height
    
    # 复制追加 PPT 的每一页
    for i, slide in enumerate(add_prs.slides, 1):
        print(f"   复制第 {i} 页...")
        
        # 使用相同的布局创建新幻灯片
        blank_layout = base_prs.slide_layouts[6]  # 空白布局
        new_slide = base_prs.slides.add_slide(blank_layout)
        
        # 复制所有形状
        for shape in slide.shapes:
            # 计算缩放比例
            add_width = add_prs.slide_width
            add_height = add_prs.slide_height
            
            scale_x = base_width / add_width if add_width != base_width else 1
            scale_y = base_height / add_height if add_height != base_height else 1
            
            # 复制形状
            if shape.shape_type == 13:  # 图片
                try:
                    image = shape.image
                    new_slide.shapes.add_picture(
                        image.blob,
                        shape.left * scale_x,
                        shape.top * scale_y,
                        shape.width * scale_x,
                        shape.height * scale_y
                    )
                except Exception as e:
                    print(f"      ⚠️ 图片复制失败: {e}")
            elif hasattr(shape, "text_frame"):
                # 文本框
                txBox = new_slide.shapes.add_textbox(
                    shape.left * scale_x,
                    shape.top * scale_y,
                    shape.width * scale_x,
                    shape.height * scale_y
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                
                # 复制文本内容
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    if para_idx == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = para.text
                    p.font.size = para.font.size
                    p.font.bold = para.font.bold
                    p.font.italic = para.font.italic
                    p.alignment = para.alignment
    
    # 保存合并后的 PPT
    base_prs.save(output_path)
    print(f"\n✅ 合并完成！")
    print(f"📊 输出文件: {output_path}")
    print(f"   总页数: {len(base_prs.slides)}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("用法: python merge_ppt.py 基础.pptx 追加.pptx [--output 输出.pptx]")
        sys.exit(1)
    
    base_ppt = sys.argv[1]
    add_ppt = sys.argv[2]
    
    if "--output" in sys.argv:
        output_idx = sys.argv.index("--output") + 1
        output_ppt = sys.argv[output_idx]
    else:
        # 默认输出文件名
        base_name = os.path.splitext(os.path.basename(base_ppt))[0]
        output_ppt = f"{base_name}_merged.pptx"
    
    merge_ppts(base_ppt, add_ppt, output_ppt)
