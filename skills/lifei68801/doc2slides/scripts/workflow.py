#!/usr/bin/env python3
# Part of doc2slides skill.
# Security: Only calls local Chrome/Chromium for HTML rendering.

"""
Enhanced workflow: Document → AI Analysis → LLM HTML Slides → PPTX

V3 Features:
- LLM-based HTML generation (not template filling)
- Support 18 layouts (Dashboard, Timeline, Comparison, etc.)
- Concurrent generation
- High-quality rendering (Playwright 3x)

Usage:
  python workflow.py --input document.pdf --output presentation.pptx
  python workflow.py --input document.pdf --output slides.pptx --preview
  python workflow.py --input document.pdf --output slides.pptx --instruction "商务风格，重点突出财务数据，控制在10页以内"
"""

import sys
import json
import argparse
import os
import asyncio
import subprocess
from pathlib import Path
from datetime import datetime

SCRIPTS_DIR = Path(__file__).parent

MIN_KEY_POINTS = 4  # 每页最少 key_points 数量

# Import LLM HTML Generator
try:
    from llm_generate_html import LLMHTMLGenerator
    HAS_LLM_GEN = True
except ImportError:
    HAS_LLM_GEN = False
    print("Warning: llm_generate_html not found, falling back to template mode")


def run_script(script_name: str, args: list, cwd=None) -> tuple:
    """Run a Python script and return (success, output)."""
    script_path = SCRIPTS_DIR / script_name
    cmd = [sys.executable, str(script_path)] + args
    
    print(f"→ {' '.join(cmd)}")
    result = subprocess.run(cmd, cwd=cwd, capture_output=True, text=True)
    
    if result.returncode != 0:
        print(f"Error: {result.stderr}")
        return False, result.stderr
    
    return True, result.stdout


def enrich_slide_data(analysis: dict) -> dict:
    """Step 2.1: Ensure every slide has enough key_points.
    
    Problem: LLM analysis sometimes generates slides with 0-2 key_points,
    causing empty/thin HTML pages. This function:
    1. Checks each slide's key_points count
    2. If < MIN_KEY_POINTS, expands from content_detail or title
    3. Also ensures data_points exist (extracts from text if missing)
    
    Returns: enriched analysis dict (modified in-place)
    """
    import re
    
    slides = analysis.get('slide_structure', analysis.get('slides', []))
    enriched_count = 0
    
    for i, slide in enumerate(slides):
        key_points = slide.get('key_points', [])
        content_detail = slide.get('content_detail', '')
        title = slide.get('title', '')
        data_points = slide.get('data_points', [])
        
        # --- Fix key_points ---
        if len(key_points) < MIN_KEY_POINTS:
            original_count = len(key_points)
            
            # Strategy 1: Split content_detail into sentences
            if content_detail:
                # Split by Chinese/Japanese punctuation, or newlines
                sentences = re.split(r'[。\n；!！？\?]', content_detail)
                sentences = [s.strip() for s in sentences if len(s.strip()) > 8]
                
                for sent in sentences:
                    if len(key_points) >= MIN_KEY_POINTS:
                        break
                    # Dedup: skip if already in key_points
                    sent_short = sent[:60]  # Truncate long sentences
                    if sent_short not in key_points and sent_short[:20] not in str(key_points):
                        key_points.append(sent_short)
            
            # Strategy 2: If title is multi-part (contains separator), split it
            if len(key_points) < MIN_KEY_POINTS and any(sep in title for sep in ['：', ':', '—', '-', '|', '/']):
                for sep in ['：', ':', '—', '-', '|', '/']:
                    if sep in title:
                        parts = title.split(sep)
                        for part in parts:
                            part = part.strip()
                            if len(part) > 4 and part not in key_points:
                                key_points.append(part)
                                if len(key_points) >= MIN_KEY_POINTS:
                                    break
                        break
            
            # Strategy 3: Extract from data_points labels
            if len(key_points) < MIN_KEY_POINTS and data_points:
                for dp in data_points:
                    label = dp.get('label', '')
                    value = dp.get('value', '')
                    unit = dp.get('unit', '')
                    point = f"{label}：{value}{unit}"
                    if len(label) > 2 and point not in key_points:
                        key_points.append(point)
                        if len(key_points) >= MIN_KEY_POINTS:
                            break
            
            # Strategy 4: Generate placeholder from slide type/title context
            slide_type = slide.get('type', slide.get('layout_suggestion', ''))
            while len(key_points) < MIN_KEY_POINTS:
                idx = len(key_points) + 1
                if slide_type in ('cover',):
                    # Cover slides should stay minimal — just ensure subtitle
                    key_points.append(f"{title}")
                    break  # Don't pad cover slides with fake points
                else:
                    key_points.append(f"{title} - 详细内容 {idx}")
            
            slide['key_points'] = key_points
            enriched_count += 1
            print(f"  + Slide {i+1} \"{title[:30]}\": key_points {original_count} → {len(key_points)}")
        
        # --- Fix data_points (extract from text if empty) ---
        if not data_points:
            combined_text = f"{content_detail} {' '.join(key_points)}"
            extracted = _extract_data_points_from_text(combined_text)
            if extracted:
                slide['data_points'] = extracted
                print(f"  + Slide {i+1} \"{title[:30]}\": extracted {len(extracted)} data_points")
    
    # Update both keys
    if 'slide_structure' in analysis:
        analysis['slide_structure'] = slides
    if 'slides' in analysis:
        analysis['slides'] = slides
    
    if enriched_count > 0:
        print(f"\n{'='*60}")
        print(f"Step 2.1: Enriched {enriched_count}/{len(slides)} slides (minimum {MIN_KEY_POINTS} key_points)")
    
    return analysis


def _extract_data_points_from_text(text: str) -> list:
    """Extract numeric data points from text using regex patterns."""
    import re
    
    if not text:
        return []
    
    data_points = []
    seen = set()
    
    patterns = [
        (r'(\d+(?:\.\d+)?)\s*(亿|万亿)(?:\s*(?:元|人民币|美元))?', '金额'),
        (r'(\d+(?:\.\d+)?)\s*万(?:\s*(?:元|人民币|美元|人|家|户|个))?', '数量'),
        (r'(\d+(?:\.\d+)?)\s*%', '占比'),
        (r'(\d+)\s*(家|人|个|名|台|套|项|款|轮|期|次)', '数量'),
        (r'(20\d{2})\s*年', '年份'),
        (r'(\d+)\s*\+', '数量'),
    ]
    
    label_keywords = {
        '亿': ['营收', '融资', '估值', '市值', '规模'],
        '万': ['营收', '客户', '用户', '收入'],
        '%': ['增长率', '占比', '覆盖率', '满意度'],
        '家': ['客户', '企业', '门店', '机构'],
        '人': ['团队', '员工', '用户', '参与'],
        '年': ['成立', '发布', '上线'],
    }
    
    for pattern, default_label in patterns:
        for match in re.finditer(pattern, text):
            value = match.group(1)
            unit = match.group(2) if len(match.groups()) > 1 else ''
            
            # Infer label from nearby context
            context_start = max(0, match.start() - 40)
            context = text[context_start:match.end() + 20]
            label = default_label
            
            for unit_key, suggestions in label_keywords.items():
                if unit_key in unit or unit_key in value:
                    for kw in suggestions:
                        if kw in context:
                            label = kw
                            break
                    break
            
            key = f"{label}:{value}{unit}"
            if key not in seen:
                seen.add(key)
                data_points.append({'label': label, 'value': value, 'unit': unit})
    
    return data_points[:6]


def read_content(input_file: str, output_file: str) -> dict:
    """Step 1: Read content from document."""
    print(f"\n{'='*60}")
    print("Step 1: Reading content...")
    
    success, output = run_script('read_content.py', [input_file, '--summarize', '--json'])
    
    if not success:
        return None
    
    # Save output to file
    try:
        content = json.loads(output)
    except json.JSONDecodeError:
        print("Failed to parse extracted content")
        return None
    
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(content, f, ensure_ascii=False, indent=2)
    
    return content


def extract_style(style_file: str, output_file: str) -> dict:
    """Step 2.5: Extract style from reference PPT (optional)."""
    print(f"\n{'='*60}")
    print("Step 2.5: Extracting style from reference PPT...")
    
    if not style_file or not os.path.exists(style_file):
        print("No style file provided, using default McKinsey style")
        return None
    
    success, _ = run_script('extract_style.py', ['--input', style_file, '--output', output_file])
    
    if not success:
        return None
    
    with open(output_file, 'r', encoding='utf-8') as f:
        return json.load(f)


async def generate_html_slides_llm(
    analysis: dict, 
    output_dir: str, 
    model: str = None,
    style: str = None,
    theme: str = None,
    instruction: str = None
) -> list:
    """Step 3: Generate HTML slides using LLM."""
    print(f"\n{'='*60}")
    print("Step 3: Generating HTML slides with LLM...")
    
    if not HAS_LLM_GEN:
        print("Error: LLM HTML Generator not available")
        return []
    
    html_dir = Path(output_dir) / "html"
    html_dir.mkdir(parents=True, exist_ok=True)
    
    slides = analysis.get('slide_structure', analysis.get('slides', []))
    html_files = []
    
    generator = LLMHTMLGenerator(model=model, style=style, theme=theme, instruction=instruction)
    
    for i, slide in enumerate(slides, 1):
        output_file = html_dir / f"slide_{i:02d}.html"
        
        # Prepare slide data
        slide_data = {
            'layout': slide.get('template', slide.get('layout_suggestion', 'CONTENT')),
            'title': slide.get('title', ''),
            'content': slide.get('key_points', []),
            'stats': slide.get('data_points', slide.get('stats', [])),
            **slide
        }
        
        # Generate HTML using LLM
        try:
            html_content = await generator.generate_slide_html(slide_data)
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(html_content)
            
            html_files.append(str(output_file))
            print(f"  ✓ Generated {output_file.name}")
        
        except Exception as e:
            print(f"  ✗ Failed to generate slide {i}: {e}")
    
    print(f"✓ Generated {len(html_files)} HTML slides with LLM")
    return html_files


def generate_html_slides(analysis: dict, style: dict, output_dir: str) -> list:
    """Step 3: Generate HTML slides (fallback to template mode)."""
    print(f"\n{'='*60}")
    print("Step 3: Generating HTML slides (template mode)...")
    
    html_dir = Path(output_dir) / "html"
    html_dir.mkdir(parents=True, exist_ok=True)
    
    slides = analysis.get('slide_structure', analysis.get('slides', []))
    html_files = []
    
    for i, slide in enumerate(slides, 1):
        output_file = html_dir / f"slide_{i:02d}.html"
        
        # Prepare slide data
        stats = slide.get('data_points', slide.get('stats', []))
        
        slide_data = {
            'template': slide.get('template', slide.get('layout_suggestion', 'CONTENT')),
            'title': slide.get('title', ''),
            'content': slide.get('key_points', []),
            'stats': stats,
            **slide
        }
        
        # Apply style if available
        if style:
            slide_data['style'] = style
        
        # Generate HTML
        success, _ = run_script('generate_html.py', [
            '--template', slide_data['template'],
            '--output', str(output_file),
            '--data', json.dumps(slide_data, ensure_ascii=False)
        ])
        
        if success:
            html_files.append(str(output_file))
    
    print(f"✓ Generated {len(html_files)} HTML slides")
    return html_files


def render_html_to_png(html_files: list, output_dir: str) -> list:
    """Render HTML files to PNG images using Playwright Python API."""
    import re
    
    preview_dir = Path(output_dir)
    preview_dir.mkdir(parents=True, exist_ok=True)
    
    png_files = []
    
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print("  ⚠ Playwright not installed, cannot render PNG")
        return []
    
    with sync_playwright() as p:
        # Launch browser with GPU acceleration disabled for stability
        browser = p.chromium.launch(args=['--disable-gpu', '--no-sandbox'])
        
        for html_file in html_files:
            html_path = Path(html_file)
            png_file = preview_dir / f"{html_path.stem}.png"
            
            # Extract dimensions from HTML
            with open(html_file, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            width_match = re.search(r'width:\s*(\d+)px', html_content)
            height_match = re.search(r'height:\s*(\d+)px', html_content)
            
            width = int(width_match.group(1)) if width_match else 1920
            height = int(height_match.group(1)) if height_match else 1080
            
            try:
                page = browser.new_page(viewport={'width': width, 'height': height})
                page.goto(f'file://{html_path.resolve()}', wait_until='networkidle')
                # Extra wait for complex SVG/CSS rendering
                page.wait_for_timeout(500)
                page.screenshot(path=str(png_file), full_page=False)
                page.close()
                png_files.append(str(png_file))
            except Exception as e:
                print(f"  ⚠ Failed to render {html_path.name}: {e}")
        
        browser.close()
    
    return png_files


def render_previews(html_files: list, output_dir: str) -> list:
    """Step 4: Render HTML to PNG previews."""
    print(f"\n{'='*60}")
    print("Step 4: Rendering previews...")
    
    preview_dir = Path(output_dir) / "preview"
    preview_dir.mkdir(parents=True, exist_ok=True)
    
    script = SCRIPTS_DIR / "html2png.sh"
    png_files = []
    
    for html_file in html_files:
        html_path = Path(html_file)
        png_file = preview_dir / f"{html_path.stem}.png"
        
        cmd = ['bash', str(script), str(html_file), str(png_file)]
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode == 0:
            png_files.append(str(png_file))
    
    if png_files:
        print(f"✓ Rendered {len(png_files)} preview images")
    else:
        print("⚠ Preview rendering skipped (Chrome not available)")
    
    return png_files


def generate_pptx_from_html(html_files: list, output_file: str, style: dict, slides_data: list) -> str:
    """Step 5: Generate PPTX by embedding rendered HTML images."""
    print(f"\n{'='*60}")
    print("Step 5: Generating PPTX...")
    
    # Try to render HTML to images
    png_files = []
    preview_dir = Path(html_files[0]).parent.parent / "preview"
    preview_dir.mkdir(parents=True, exist_ok=True)
    
    # Check if PNGs already exist
    for html_file in html_files:
        html_path = Path(html_file)
        png_file = preview_dir / f"{html_path.stem}.png"
        if png_file.exists():
            png_files.append(str(png_file))
    
    # If not all PNGs exist, render them
    if len(png_files) != len(html_files):
        print("  Rendering HTML to PNG images...")
        png_files = render_html_to_png(html_files, str(preview_dir))
    
    if png_files and len(png_files) == len(html_files):
        # Use PNG embedding approach (true HTML fusion)
        return generate_pptx_from_pngs(png_files, output_file, style)
    else:
        # Fallback to native PPTX generation
        print("  ⚠ PNG rendering failed, using native PPTX generation")
        return generate_pptx_native(slides_data, output_file, style)


def generate_pptx_from_pngs(png_files: list, output_file: str, style: dict) -> str:
    """Generate PPTX by embedding PNG images (true HTML fusion)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("Error: python-pptx not installed")
        return None
    
    # Create presentation with correct size
    if style and 'layout' in style:
        slide_width = Inches(style['layout'].get('slide_size', {}).get('width', 13.33))
        slide_height = Inches(style['layout'].get('slide_size', {}).get('height', 7.5))
    else:
        slide_width = Inches(13.33)
        slide_height = Inches(7.5)
    
    prs = Presentation()
    prs.slide_width = slide_width
    prs.slide_height = slide_height
    
    # Add slides with PNG images
    blank_layout = prs.slide_layouts[6]  # Blank layout
    
    for png_file in png_files:
        slide = prs.slides.add_slide(blank_layout)
        
        # Add PNG as full-slide background
        slide.shapes.add_picture(png_file, 0, 0, slide_width, slide_height)
    
    prs.save(output_file)
    print(f"✓ Generated PPTX with {len(png_files)} slides (HTML fusion)")
    return output_file


def generate_pptx_native(slides_data: list, output_file: str, style: dict) -> str:
    """Generate PPTX using native generator (fallback)."""
    # Prepare slides JSON
    slides_json = {'slides': slides_data, 'style': style}
    slides_file = Path(output_file).parent / "slides_data.json"
    
    with open(slides_file, 'w', encoding='utf-8') as f:
        json.dump(slides_json, f, ensure_ascii=False, indent=2)
    
    success, _ = run_script('generate_pptx_pro.py', [
        '--output', output_file,
        '--content', str(slides_file)
    ])
    
    if success:
        print(f"✓ Generated PPTX (native generation)")
        return output_file
    
    return None


def validate_pptx(pptx_file: str) -> dict:
    """Step 6: Validate PPTX quality."""
    print(f"\n{'='*60}")
    print("Step 6: Validating PPTX quality...")
    
    report_file = Path(pptx_file).parent / "validation_report.json"
    
    success, output = run_script('validate.py', ['--input', pptx_file, '--output', str(report_file)])
    
    if not success:
        return {'passed': False, 'error': output}
    
    with open(report_file, 'r', encoding='utf-8') as f:
        report = json.load(f)
    
    # Print summary
    if report.get('passed'):
        print(f"✅ PPT 质量验证通过！评分: {report.get('score', 0)}/100")
    else:
        print(f"⚠️ 发现 {report.get('issues_count', 0)} 个问题：")
        for issue in report.get('issues', [])[:3]:
            print(f"  - {issue['message']}")
    
    return report


def main():
    parser = argparse.ArgumentParser(description='Enhanced document to PPT workflow')
    parser.add_argument('--input', required=True, help='Input file (PDF/Word/PPT/JSON)')
    parser.add_argument('--output', required=True, help='Output PPTX file')
    parser.add_argument('--style', help='Reference PPTX file for style extraction OR color scheme name (corporate/tech/nature/warm/minimal/dark_purple/finance)')
    parser.add_argument('--theme', help='Document theme for auto color selection (e.g., AI, finance, health)')
    parser.add_argument('--model', help='AI model for content analysis')
    parser.add_argument('--preview', action='store_true', help='Generate PNG previews')
    parser.add_argument('--skip-validate', action='store_true', help='Skip quality validation')
    parser.add_argument('--instruction', help='User instruction to guide the entire PPT generation pipeline. This overrides default behavior in content analysis, layout selection, and HTML generation. Example: "商务风格，重点突出财务数据，控制在10页以内，每页必须有图表"')
    
    args = parser.parse_args()
    
    # Setup directories
    output_dir = Path(args.output).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    
    work_dir = output_dir / "work"
    work_dir.mkdir(parents=True, exist_ok=True)
    
    content_file = work_dir / "content.json"
    analysis_file = work_dir / "analysis.json"
    style_file = work_dir / "style.json"
    html_dir = work_dir / "slides"
    
    # Step 1: Extract content or load slides directly
    input_path = Path(args.input)
    analysis = None
    
    # If input is JSON with slides, load directly
    if input_path.suffix == '.json':
        print(f"\n{'='*60}")
        print("Step 1: Loading slides from JSON...")
        with open(args.input, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Support both {slides: [...]} and {slide_structure: [...]}
        analysis = {
            'slides': data.get('slides', data.get('slide_structure', [])),
            'slide_structure': data.get('slides', data.get('slide_structure', []))
        }
        print(f"  ✓ Loaded {len(analysis['slides'])} slides from JSON")
    
    else:
        # Read from document (PDF/Word/PPT)
        content = read_content(args.input, str(content_file))
        if not content:
            print("Failed to read content")
            sys.exit(1)
    
    # Step 2: AI analysis (only for non-JSON inputs)
    if not analysis:
        # Check for existing outline.json first
        outline_candidates = [
            input_path.parent / "outline.json",
            input_path.parent.parent / "outline.json",
            Path(args.output).parent / "outline.json",
            Path("output/outline.json"),
        ]
        
        for outline_path in outline_candidates:
            if outline_path.exists():
                print(f"  Found existing outline: {outline_path}")
                with open(outline_path, 'r', encoding='utf-8') as f:
                    analysis = json.load(f)
                break
        
        if not analysis:
            # Step 2: Use LLM to analyze content and generate outline
            print(f"\n{'='*60}")
            print("Step 2: Analyzing content with LLM...")
            
            outline_file = Path("output/outline.json")
            outline_file.parent.mkdir(parents=True, exist_ok=True)
            
            # Save content to temp file for analysis
            content_temp = Path("output/content_temp.json")
            with open(content_temp, 'w', encoding='utf-8') as f:
                json.dump(content, f, ensure_ascii=False)
            
            analyze_args = [
                '--content-json', str(content_temp),
                '--output', str(outline_file),
                '--model', args.model or ''
            ]
            if args.instruction:
                analyze_args.extend(['--instruction', args.instruction])
            
            success, _ = run_script('analyze_content.py', analyze_args)
            
            if success and outline_file.exists():
                with open(outline_file, 'r', encoding='utf-8') as f:
                    analysis = json.load(f)
                print(f"  ✓ Generated outline with {len(analysis.get('slides', []))} slides")
            else:
                # Fallback: simple structure
                print("  ⚠ LLM analysis failed, using simple structure")
                text = content.get('text', '')
                slides = []
                if text:
                    slides.append({
                        'title': content.get('title', '内容概览'),
                        'template': 'CONTENT',
                        'key_points': text[:500].split('\n')[:5]
                    })
                analysis = {'slides': slides, 'slide_structure': slides}
    
    # Step 2.1: Enrich slide data (ensure every slide has enough key_points)
    analysis = enrich_slide_data(analysis)
    
    # Save enriched analysis for debugging
    if analysis_file:
        with open(str(analysis_file), 'w', encoding='utf-8') as f:
            json.dump(analysis, f, ensure_ascii=False, indent=2)
    
    # Step 2.5: Extract style (optional)
    color_style = None  # 配色方案名称
    ref_style = None    # 参考 PPT 提取的样式
    
    if args.style:
        # 判断是配色方案名称还是 PPT 文件路径
        if args.style.endswith('.pptx') or args.style.endswith('.ppt'):
            # 是 PPT 文件路径，提取样式
            ref_style = extract_style(args.style, str(style_file))
        else:
            # 是配色方案名称
            color_style = args.style
    
    # Step 3: Generate HTML (LLM mode or fallback to template)
    if HAS_LLM_GEN:
        html_files = asyncio.run(generate_html_slides_llm(
            analysis, str(html_dir), args.model, style=color_style, theme=args.theme,
            instruction=args.instruction
        ))
    else:
        html_files = generate_html_slides(analysis, ref_style, str(html_dir))
    
    if not html_files:
        print("Failed to generate HTML slides")
        sys.exit(1)
    
    # Step 4: Render previews (optional)
    if args.preview:
        render_previews(html_files, str(html_dir))
    
    # Step 5: Generate PPTX
    slides_data = analysis.get('slide_structure', analysis.get('slides', []))
    # If output is a directory, create a filename
    output_path = Path(args.output)
    if output_path.is_dir():
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_output = str(output_path / f"slides_{timestamp}.pptx")
    else:
        pptx_output = args.output
    
    # Use the style dict (color_style or ref_style)
    pptx_style = ref_style if ref_style else {'theme': color_style or 'default'}
    pptx_file = generate_pptx_from_html(html_files, pptx_output, pptx_style, slides_data)
    
    if not pptx_file:
        print("Failed to generate PPTX")
        sys.exit(1)
    
    # Step 6: Validate (optional)
    if not args.skip_validate:
        validate_pptx(pptx_file)
    
    # Summary
    print(f"\n{'='*60}")
    print(f"✓ Complete! Output: {pptx_file}")
    print(f"  HTML slides: {html_dir / 'html'}")
    print(f"  Work files: {work_dir}")


if __name__ == "__main__":
    main()
