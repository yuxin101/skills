#!/usr/bin/env python3
# Part of doc2slides skill.

#!/usr/bin/env python3
"""
png2pptx.py - Create PPTX from PNG images (one PNG per slide)
PNG images are used as slide backgrounds
"""

import argparse
import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Standard 16:9 slide dimensions (in EMUs)
SLIDE_WIDTH_EMU = 12192000  # 13.333 inches
SLIDE_HEIGHT_EMU = 6858000  # 7.5 inches
SLIDE_WIDTH_INCH = 13.333
SLIDE_HEIGHT_INCH = 7.5


def create_presentation():
    """Create a new presentation with 16:9 aspect ratio"""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
    return prs


def add_slide_with_background(prs: Presentation, png_path: str, slide_num: int = None):
    """
    Add a slide with PNG image as background
    
    Args:
        prs: Presentation object
        png_path: Path to PNG image
        slide_num: Optional slide number (for reference)
    """
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Add image as full-slide background
    # Position: (left, top, width, height) in EMUs
    slide.shapes.add_picture(
        png_path,
        Emu(0),
        Emu(0),
        Emu(SLIDE_WIDTH_EMU),
        Emu(SLIDE_HEIGHT_EMU)
    )
    
    return slide


def create_pptx_from_pngs(png_dir: str, output_path: str, metadata: dict = None):
    """
    Create PPTX from all PNG images in a directory
    
    Args:
        png_dir: Directory containing PNG files (slide_01.png, slide_02.png, etc.)
        output_path: Output PPTX file path
        metadata: Optional metadata (title, author, etc.)
    """
    png_dir = Path(png_dir)
    
    # Find all PNG files, sorted by name
    png_files = sorted(png_dir.glob('*.png'))
    
    if not png_files:
        raise ValueError(f"No PNG files found in {png_dir}")
    
    print(f"\n=== Creating PPTX from {len(png_files)} PNG slides ===")
    print(f"  Input: {png_dir}")
    print(f"  Output: {output_path}")
    
    # Create presentation
    prs = create_presentation()
    
    # Add metadata if provided
    if metadata:
        if 'title' in metadata:
            prs.core_properties.title = metadata['title']
        if 'author' in metadata:
            prs.core_properties.author = metadata['author']
        if 'subject' in metadata:
            prs.core_properties.subject = metadata['subject']
    
    # Add each PNG as a slide
    for i, png_file in enumerate(png_files, 1):
        print(f"  [{i}/{len(png_files)}] {png_file.name}")
        add_slide_with_background(prs, str(png_file), slide_num=i)
    
    # Save
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    
    size_kb = os.path.getsize(output_path) / 1024
    print(f"\n✓ Created: {output_path} ({size_kb:.1f}KB)")
    
    return {
        'success': True,
        'output': str(output_path),
        'total_slides': len(png_files),
        'file_size_kb': size_kb
    }


def create_pptx_from_json(png_dir: str, output_path: str, slides_json: str = None):
    """
    Create PPTX from PNG images with optional slide metadata
    
    Args:
        png_dir: Directory containing PNG files
        output_path: Output PPTX file path
        slides_json: Optional JSON file with slide metadata
    """
    metadata = None
    
    if slides_json and os.path.exists(slides_json):
        with open(slides_json, 'r', encoding='utf-8') as f:
            data = json.load(f)
            metadata = data.get('metadata', {})
    
    return create_pptx_from_pngs(png_dir, output_path, metadata)


def main():
    parser = argparse.ArgumentParser(description='Create PPTX from PNG slide images')
    parser.add_argument('png_dir', help='Directory containing PNG files')
    parser.add_argument('--output', '-o', required=True, help='Output PPTX file path')
    parser.add_argument('--slides', help='Optional JSON file with slide metadata')
    
    args = parser.parse_args()
    
    result = create_pptx_from_json(args.png_dir, args.output, args.slides)
    print(f"\n{json.dumps(result, ensure_ascii=False)}")


if __name__ == '__main__':
    main()
