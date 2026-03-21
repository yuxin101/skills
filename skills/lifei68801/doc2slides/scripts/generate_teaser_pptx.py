#!/usr/bin/env python3
# Part of doc2slides skill.

#!/usr/bin/env python3
"""
Generate Teaser-style PowerPoint presentation.
Dark gradient background + large title + minimal content.
Direct PPTX output, no HTML intermediate.

Usage:
  python generate_teaser_pptx.py --output slides.pptx --content slides.json
  python generate_teaser_pptx.py --output slides.pptx --slides '[{"title":"...", "content":[...]}]'
"""

import sys
import json
import argparse
import os

def create_teaser_presentation(slides, output_path="output.pptx", title="数势科技"):
    """
    Create Teaser-style PowerPoint with dark background.
    
    Args:
        slides: List of slide dicts with 'title', 'content'
        output_path: Output PPTX file path
        title: Presentation title
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        return {"success": False, "error": "python-pptx not installed. Run: pip install python-pptx"}
    
    try:
        # Create presentation with 16:9 aspect ratio
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Colors (Teaser style - dark navy)
        NAVY_900 = RGBColor(0x0f, 0x17, 0x2a)   # Very dark navy
        NAVY_800 = RGBColor(0x1e, 0x29, 0x3b)   # Dark navy
        BLUE_500 = RGBColor(0x3b, 0x82, 0xf6)   # Bright blue accent
        WHITE = RGBColor(0xff, 0xff, 0xff)
        SLATE_300 = RGBColor(0xcb, 0xd5, 0xe1)  # Light gray
        
        for idx, slide_data in enumerate(slides):
            slide_title = slide_data.get('title', f'Slide {idx + 1}')
            slide_content = slide_data.get('content', [])
            
            # Use blank layout
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # === Background: Dark navy gradient (simulated with solid color) ===
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = NAVY_900
            
            # === Page number (top right) ===
            page_box = slide.shapes.add_textbox(
                Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.8)
            )
            tf = page_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{idx + 1:02d}"
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x33, 0x47, 0x63)  # Darker navy for subtle look
            p.alignment = PP_ALIGN.RIGHT
            
            # === Eyebrow label ===
            eyebrow_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.5), Inches(6), Inches(0.5)
            )
            tf = eyebrow_box.text_frame
            p = tf.paragraphs[0]
            p.text = "数势科技 · DigitForce"
            p.font.size = Pt(12)
            p.font.color.rgb = BLUE_500
            p.font.bold = True
            
            # === Main title ===
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.2), Inches(11), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.line_spacing = 1.2
            
            # === Content bullets ===
            if slide_content:
                content_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(4.2), Inches(10), Inches(2.5)
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                
                # Handle string or list content
                if isinstance(slide_content, str):
                    items = [line.strip() for line in slide_content.split('\n') if line.strip()]
                else:
                    items = slide_content
                
                for i, item in enumerate(items[:5]):  # Max 5 items
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    # Bullet point
                    p.text = f"●  {item}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = SLATE_300
                    p.space_after = Pt(12)
                    p.line_spacing = 1.3
            
            # === Bottom accent line ===
            accent_line = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.8), Inches(6.5), Inches(2), Inches(0.05)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = BLUE_500
            accent_line.line.fill.background()  # No border
            
            # === Bottom text ===
            bottom_box = slide.shapes.add_textbox(
                Inches(3), Inches(6.4), Inches(5), Inches(0.4)
            )
            tf = bottom_box.text_frame
            p = tf.paragraphs[0]
            p.text = "AI Agent · 数据智能决策"
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(0x64, 0x74, 0x8b)  # Slate-500
        
        prs.save(output_path)
        
        return {
            "success": True,
            "output": output_path,
            "total_slides": len(prs.slides),
            "message": f"Created {len(prs.slides)} Teaser-style slides"
        }
        
    except Exception as e:
        import traceback
        return {"success": False, "error": str(e), "traceback": traceback.format_exc()}

def main():
    parser = argparse.ArgumentParser(description='Generate Teaser-style PowerPoint')
    parser.add_argument('--output', '-o', required=True, help='Output PPTX file path')
    parser.add_argument('--content', help='JSON file with slides array')
    parser.add_argument('--slides', help='JSON string with slides array')
    parser.add_argument('--title', default='数势科技', help='Presentation title')
    
    args = parser.parse_args()
    
    # Get slides
    slides = []
    
    if args.slides:
        if os.path.exists(args.slides):
            with open(args.slides, 'r', encoding='utf-8') as f:
                slides = json.load(f)
        else:
            try:
                slides = json.loads(args.slides)
            except:
                slides = [{"title": "Content", "content": args.slides}]
    elif args.content:
        with open(args.content, 'r', encoding='utf-8') as f:
            data = json.load(f)
            # Handle both list and dict formats
            if isinstance(data, list):
                slides = data
            else:
                slides = data.get('slides', [data])
    else:
        # Read from stdin
        try:
            data = json.load(sys.stdin)
            if isinstance(data, list):
                slides = data
            else:
                slides = data.get('slides', [data])
        except:
            print(json.dumps({"success": False, "error": "No content provided"}))
            sys.exit(1)
    
    result = create_teaser_presentation(
        slides=slides,
        output_path=args.output,
        title=args.title
    )
    
    print(json.dumps(result, ensure_ascii=False, indent=2))

if __name__ == "__main__":
    main()
