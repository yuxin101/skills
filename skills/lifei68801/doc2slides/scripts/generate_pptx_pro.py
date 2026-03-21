#!/usr/bin/env python3
# Part of doc2slides skill.

#!/usr/bin/env python3
"""
Professional PPTX Generator - McKinsey/BCG Style
Converts document content to polished consulting-style PowerPoint.

Templates:
1. COVER - Title slide with company branding
2. PROBLEM - Data-driven problem statement (left: text, right: big numbers + progress bars)
3. SOLUTION - 5-step solution flow with benefits
4. CONTENT - Standard bullet point content
5. MATRIX - 2x2 strategic matrix
6. COMPARE - Before/After or Do/Don't comparison
7. CLOSING - Call to action / contact slide

Usage:
  python generate_pptx_pro.py --output slides.pptx --content slides.json
  python generate_pptx_pro.py --output slides.pptx --slides '[{"title":"...", "content":[...]}]'
"""

import sys
import json
import argparse
import os
import re

# Color Palette - McKinsey/BCG Style
COLORS = {
    'navy_900': '0f172a',
    'navy_800': '1e293b',
    'navy_700': '334155',
    'slate_600': '475569',
    'slate_400': '94a3b8',
    'slate_200': 'e2e8f0',
    'slate_100': 'f1f5f9',
    'white': 'ffffff',
    'red': 'dc2626',
    'orange': 'ea580c',
    'amber': 'd97706',
    'green': '059669',
    'emerald': '10b981',
    'blue': '3b82f6',
}


def rgb(hex_color):
    """Convert hex color to RGB tuple for python-pptx"""
    from pptx.dml.color import RGBColor
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def add_shape_with_color(slide, shape_type, left, top, width, height, fill_color, line_color=None):
    """Add a shape with specified fill color"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
    else:
        shape.line.fill.background()  # No border
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, font_color='475569', bold=False, align='left', font_name='Microsoft YaHei'):
    """Add a text box with specified formatting"""
    from pptx.enum.text import PP_ALIGN
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = rgb(font_color)
    p.font.bold = bold
    p.font.name = font_name
    
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    
    return txBox


def detect_template(slide_data):
    """Intelligently detect the best template for content"""
    # 优先使用数据中已有的 template 字段
    if slide_data.get('template'):
        return slide_data['template'].upper()
    
    title = slide_data.get('title', '').lower()
    content = slide_data.get('content', []) or slide_data.get('key_points', [])
    
    if isinstance(content, str):
        content = [content]
    
    content_str = ' '.join(content).lower()
    
    # Check for numbers/stats
    has_numbers = bool(re.search(r'\d+[%％]', content_str))
    has_stats = any(re.search(r'\d+', c) for c in content[:3]) if content else False
    
    # Detect template type
    if any(kw in title for kw in ['封面', 'cover', '介绍', 'introduction']):
        return 'COVER'
    elif any(kw in title for kw in ['痛点', '问题', '挑战', 'problem', '挑战']):
        return 'PROBLEM'
    elif any(kw in title for kw in ['方案', 'solution', '解决方案', '路径']):
        return 'SOLUTION'
    elif any(kw in title for kw in ['对比', 'compare', 'vs', 'versus', 'do', "don't"]):
        return 'COMPARE'
    elif any(kw in title for kw in ['矩阵', 'matrix', '象限', 'quadrant']):
        return 'MATRIX'
    elif any(kw in title for kw in ['流程', 'process', '步骤', 'step']):
        return 'PROCESS'
    elif any(kw in title for kw in ['结束', 'closing', '联系', 'contact', '封底']):
        return 'CLOSING'
    elif has_numbers and has_stats:
        return 'PROBLEM'
    else:
        return 'CONTENT'


def extract_stats(content):
    """Extract statistics from content for PROBLEM template"""
    stats = []
    for item in content[:5]:
        # Look for percentage patterns
        match = re.search(r'(\d+)\s*[%％]\s*[-—]?\s*(.+)', item)
        if match:
            stats.append({
                'value': int(match.group(1)),
                'label': match.group(2).strip()[:30]
            })
        # Look for standalone numbers
        elif re.search(r'\d+', item):
            match = re.search(r'(\d+[,.]?\d*)\s*(.+)', item)
            if match:
                stats.append({
                    'value': match.group(1).replace(',', ''),
                    'label': match.group(2).strip()[:30]
                })
    return stats


def create_cover_slide(prs, slide_data, slide_num):
    """Create cover slide with company branding"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background: Navy gradient effect (solid navy)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(COLORS['navy_900'])
    bg.line.fill.background()
    
    # Company label (eyebrow)
    add_text_box(
        slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5),
        "数势科技 · DigitForce",
        font_size=Pt(14), font_color=COLORS['blue'], bold=True
    )
    
    # Main title
    title = slide_data.get('title', '')
    add_text_box(
        slide, Inches(0.8), Inches(2.5), Inches(11), Inches(2),
        title,
        font_size=Pt(48), font_color=COLORS['white'], bold=True
    )
    
    # Subtitle / Tagline
    content = slide_data.get('content', [])
    if content:
        subtitle = content[0] if isinstance(content, list) else content
        add_text_box(
            slide, Inches(0.8), Inches(4.5), Inches(11), Inches(1),
            subtitle[:80],
            font_size=Pt(24), font_color=COLORS['slate_400']
        )
    
    # Bottom accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(6.3), Inches(3), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS['blue'])
    line.line.fill.background()
    
    # Bottom text
    add_text_box(
        slide, Inches(4), Inches(6.2), Inches(8), Inches(0.5),
        "AI Agent 驱动的数据智能决策平台",
        font_size=Pt(14), font_color=COLORS['slate_400']
    )
    
    # Page number
    add_text_box(
        slide, Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.8),
        f"{slide_num:02d}",
        font_size=Pt(36), font_color=COLORS['navy_700'], bold=True, align='right'
    )
    
    return slide


def create_problem_slide(prs, slide_data, slide_num):
    """Create problem statement slide with data visualization"""
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(COLORS['white'])
    bg.line.fill.background()
    
    # === Left Column (5/12) ===
    left_col_x = Inches(0.6)
    left_col_w = Inches(5.5)
    
    # Eyebrow
    add_text_box(
        slide, left_col_x, Inches(0.6), left_col_w, Inches(0.4),
        slide_data.get('eyebrow', '数据分析'),
        font_size=Pt(12), font_color=COLORS['slate_400'], bold=True
    )
    
    # Title
    add_text_box(
        slide, left_col_x, Inches(1.0), left_col_w, Inches(1),
        slide_data.get('title', ''),
        font_size=Pt(32), font_color=COLORS['navy_900'], bold=True
    )
    
    # Description - 支持 key_points 和 content 两种字段
    content = slide_data.get('content', []) or slide_data.get('key_points', [])
    if isinstance(content, list):
        desc_text = '\n'.join([f"• {c}" for c in content[:3]])
    else:
        desc_text = content
    
    desc_box = add_text_box(
        slide, left_col_x, Inches(2.2), left_col_w, Inches(2.5),
        desc_text[:300],
        font_size=Pt(16), font_color=COLORS['slate_600']
    )
    
    # Insight box (navy background)
    insight_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left_col_x, Inches(5.2), left_col_w, Inches(1.2)
    )
    insight_bg.fill.solid()
    insight_bg.fill.fore_color.rgb = rgb(COLORS['navy_900'])
    insight_bg.line.fill.background()
    
    add_text_box(
        slide, Inches(0.8), Inches(5.35), left_col_w - Inches(0.3), Inches(0.3),
        "💡 关键洞察",
        font_size=Pt(11), font_color=COLORS['slate_400'], bold=True
    )
    
    insight_text = content[-1] if content and len(content) > 3 else "数据驱动决策是关键"
    add_text_box(
        slide, Inches(0.8), Inches(5.65), left_col_w - Inches(0.3), Inches(0.6),
        insight_text[:50],
        font_size=Pt(14), font_color=COLORS['white']
    )
    
    # === Right Column (7/12) - Stats ===
    right_col_x = Inches(6.5)
    stats = slide_data.get('stats', [])  # 优先使用 stats 字段
    if not stats:
        stats = extract_stats(content) if content else []
    
    # If no stats extracted, create placeholder
    if not stats:
        stats = [
            {'value': '73', 'label': '数据孤岛问题'},
            {'value': '65', 'label': '系统割裂严重'},
            {'value': '52', 'label': '决策效率低下'},
        ]
    
    stat_colors = [COLORS['red'], COLORS['orange'], COLORS['amber']]
    
    for i, stat in enumerate(stats[:3]):
        y_offset = Inches(1.5 + i * 1.5)
        
        # Big number
        value_text = f"{stat['value']}" if isinstance(stat['value'], str) else f"{stat['value']}"
        add_text_box(
            slide, right_col_x, y_offset, Inches(1.5), Inches(0.8),
            value_text,
            font_size=Pt(48), font_color=stat_colors[i], bold=True, align='right'
        )
        
        # Percent sign
        add_text_box(
            slide, right_col_x + Inches(1.5), y_offset + Inches(0.15), Inches(0.5), Inches(0.5),
            "%",
            font_size=Pt(24), font_color=stat_colors[i], bold=True
        )
        
        # Label
        add_text_box(
            slide, right_col_x + Inches(2.2), y_offset + Inches(0.1), Inches(4), Inches(0.4),
            stat['label'],
            font_size=Pt(18), font_color=COLORS['navy_900'], bold=True
        )
        
        # Progress bar background
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            right_col_x + Inches(2.2), y_offset + Inches(0.5), 
            Inches(4), Pt(12)
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = rgb(COLORS['slate_200'])
        bar_bg.line.fill.background()
        
        # Progress bar fill
        if isinstance(stat['value'], int):
            bar_width = Inches(4 * min(stat['value'], 100) / 100)
            bar_fill = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                right_col_x + Inches(2.2), y_offset + Inches(0.5),
                bar_width, Pt(12)
            )
            bar_fill.fill.solid()
            bar_fill.fill.fore_color.rgb = rgb(stat_colors[i])
            bar_fill.line.fill.background()
    
    # Source note
    add_text_box(
        slide, right_col_x, Inches(6.5), Inches(5.5), Inches(0.3),
        slide_data.get('source', '数据来源：IDC 2025年报告'),
        font_size=Pt(10), font_color=COLORS['slate_400']
    )
    
    # Page number
    add_text_box(
        slide, Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.8),
        f"{slide_num:02d}",
        font_size=Pt(36), font_color=COLORS['slate_200'], bold=True, align='right'
    )
    
    return slide


def create_solution_slide(prs, slide_data, slide_num):
    """Create solution flow slide with 5 steps"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(COLORS['white'])
    bg.line.fill.background()
    
    # Header
    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.4),
        slide_data.get('eyebrow', '解决方案'),
        font_size=Pt(12), font_color=COLORS['slate_400'], bold=True
    )
    
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        slide_data.get('title', ''),
        font_size=Pt(28), font_color=COLORS['navy_900'], bold=True
    )
    
    # Definition box
    def_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5), Inches(0.5)
    )
    def_bg.fill.solid()
    def_bg.fill.fore_color.rgb = rgb(COLORS['slate_100'])
    def_bg.line.color.rgb = rgb(COLORS['slate_200'])
    
    add_text_box(
        slide, Inches(1), Inches(1.88), Inches(4.5), Inches(0.3),
        "SwiftAgent = AI 数据分析师",
        font_size=Pt(12), font_color=COLORS['navy_800']
    )
    
    # 5-Step Flow - 支持 key_points 和 content 两种字段
    content = slide_data.get('content', []) or slide_data.get('key_points', [])
    if isinstance(content, str):
        content = [content]
    
    # Extract 5 steps from content or create placeholders
    steps = []
    for i, item in enumerate(content[:5]):
        steps.append({
            'number': i + 1,
            'title': item[:20] if len(item) > 20 else item,
            'detail': item[20:50] if len(item) > 20 else ''
        })
    
    # Fill remaining steps if needed
    while len(steps) < 5:
        steps.append({
            'number': len(steps) + 1,
            'title': f'步骤 {len(steps) + 1}',
            'detail': ''
        })
    
    step_colors = [COLORS['navy_900'], COLORS['navy_800'], COLORS['navy_700'], COLORS['navy_700'], COLORS['green']]
    
    for i, step in enumerate(steps):
        x = Inches(0.8 + i * 2.4)
        y = Inches(2.6)
        
        # Step card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = rgb(step_colors[i])
        card.line.fill.background()
        
        # Step number
        add_text_box(
            slide, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.4),
            f"0{i+1}",
            font_size=Pt(12), font_color=COLORS['slate_400'], bold=True
        )
        
        # Step title
        add_text_box(
            slide, x + Inches(0.2), y + Inches(0.6), Inches(1.8), Inches(0.6),
            step['title'],
            font_size=Pt(16), font_color=COLORS['white'], bold=True
        )
        
        # Step detail
        if step['detail']:
            add_text_box(
                slide, x + Inches(0.2), y + Inches(1.3), Inches(1.8), Inches(1),
                step['detail'][:40],
                font_size=Pt(11), font_color=COLORS['slate_400']
            )
        
        # Arrow between steps
        if i < 4:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + Inches(2.25), y + Inches(1.1), Inches(0.3), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = rgb(COLORS['slate_400'])
            arrow.line.fill.background()
    
    # Benefits row
    benefits = [
        {'number': '01', 'title': '效率提升', 'detail': '3-5天 → 秒级'},
        {'number': '02', 'title': '成本降低', 'detail': 'ROI 300-500%'},
        {'number': '03', 'title': '决策准确', 'detail': '数据驱动'},
    ]
    
    for i, benefit in enumerate(benefits):
        x = Inches(0.8 + i * 4.2)
        y = Inches(5.5)
        
        # Benefit card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(1.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = rgb('ecfdf5')  # emerald-50
        card.line.color.rgb = rgb('a7f3d0')  # emerald-200
        
        # Benefit number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.35), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(COLORS['emerald'])
        circle.line.fill.background()
        
        add_text_box(
            slide, x + Inches(0.22), y + Inches(0.43), Inches(0.4), Inches(0.3),
            benefit['number'],
            font_size=Pt(10), font_color=COLORS['white'], bold=True
        )
        
        # Benefit text
        add_text_box(
            slide, x + Inches(0.8), y + Inches(0.25), Inches(2.8), Inches(0.4),
            benefit['title'],
            font_size=Pt(14), font_color='065f46', bold=True  # emerald-800
        )
        
        add_text_box(
            slide, x + Inches(0.8), y + Inches(0.6), Inches(2.8), Inches(0.4),
            benefit['detail'],
            font_size=Pt(12), font_color='059669'  # emerald-600
        )
    
    # Page number
    add_text_box(
        slide, Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.8),
        f"{slide_num:02d}",
        font_size=Pt(36), font_color=COLORS['slate_200'], bold=True, align='right'
    )
    
    return slide


def create_content_slide(prs, slide_data, slide_num):
    """Create standard content slide with bullet points"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(COLORS['white'])
    bg.line.fill.background()
    
    # Header with accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(2), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS['blue'])
    line.line.fill.background()
    
    # Eyebrow
    add_text_box(
        slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
        slide_data.get('eyebrow', '数势科技'),
        font_size=Pt(12), font_color=COLORS['slate_400'], bold=True
    )
    
    # Title
    add_text_box(
        slide, Inches(0.8), Inches(1.1), Inches(11), Inches(1),
        slide_data.get('title', ''),
        font_size=Pt(32), font_color=COLORS['navy_900'], bold=True
    )
    
    # Content bullets - 支持 key_points 和 content 两种字段
    content = slide_data.get('content', []) or slide_data.get('key_points', [])
    if isinstance(content, str):
        content = [content]
    
    y_offset = Inches(2.4)
    for i, item in enumerate(content[:6]):
        # Bullet point
        bullet = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.8), y_offset + Inches(0.15), Pt(10), Pt(10)
        )
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = rgb(COLORS['blue'])
        bullet.line.fill.background()
        
        # Text
        add_text_box(
            slide, Inches(1.3), y_offset, Inches(10.5), Inches(0.8),
            item[:120],
            font_size=Pt(18), font_color=COLORS['slate_600']
        )
        
        y_offset += Inches(0.7)
    
    # Bottom decoration
    add_text_box(
        slide, Inches(0.8), Inches(6.5), Inches(6), Inches(0.3),
        "AI Agent · 数据智能决策",
        font_size=Pt(11), font_color=COLORS['slate_400']
    )
    
    # Page number
    add_text_box(
        slide, Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.8),
        f"{slide_num:02d}",
        font_size=Pt(36), font_color=COLORS['slate_200'], bold=True, align='right'
    )
    
    return slide


def create_matrix_slide(prs, slide_data, slide_num):
    """Create 2x2 matrix slide"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # White background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(COLORS['white'])
    bg.line.fill.background()
    
    # Header
    add_text_box(
        slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.4),
        slide_data.get('eyebrow', '战略矩阵'),
        font_size=Pt(12), font_color=COLORS['slate_400'], bold=True
    )
    
    add_text_box(
        slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
        slide_data.get('title', ''),
        font_size=Pt(28), font_color=COLORS['navy_900'], bold=True
    )
    
    # 2x2 Grid
    quadrants = slide_data.get('quadrants', [])
    positions = [
        {'x': Inches(0.8), 'y': Inches(2.0)},      # top-left
        {'x': Inches(6.6), 'y': Inches(2.0)},      # top-right
        {'x': Inches(0.8), 'y': Inches(4.7)},      # bottom-left
        {'x': Inches(6.6), 'y': Inches(4.7)},      # bottom-right
    ]
    
    colors = [COLORS['blue'], COLORS['green'], COLORS['orange'], COLORS['navy_800']]
    
    for i, quad in enumerate(quadrants[:4]):
        if i >= len(positions):
            break
        
        pos = positions[i]
        w = Inches(5.5)
        h = Inches(2.5)
        
        # Quadrant card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, pos['x'], pos['y'], w, h
        )
        card.fill.solid()
        card.fill.fore_color.rgb = rgb('f8fafc')  # slate-50
        card.line.color.rgb = rgb(colors[i])
        
        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, pos['x'], pos['y'], w, Inches(0.6)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = rgb(colors[i])
        title_bar.line.fill.background()
        
        # Title
        add_text_box(
            slide, pos['x'] + Inches(0.2), pos['y'] + Inches(0.15), w - Inches(0.4), Inches(0.4),
            quad.get('title', f'象限 {i+1}'),
            font_size=Pt(16), font_color=COLORS['white'], bold=True
        )
        
        # Points
        points = quad.get('points', [])
        y_offset = pos['y'] + Inches(0.7)
        for point in points[:3]:
            add_text_box(
                slide, pos['x'] + Inches(0.3), y_offset, w - Inches(0.6), Inches(0.5),
                f"• {point}",
                font_size=Pt(13), font_color=COLORS['slate_600']
            )
            y_offset += Inches(0.5)
    
    # Page number
    add_text_box(
        slide, Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.8),
        f"{slide_num:02d}",
        font_size=Pt(36), font_color=COLORS['slate_200'], bold=True, align='right'
    )
    
    return slide


def create_closing_slide(prs, slide_data, slide_num):
    """Create closing slide with call to action"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Navy background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(COLORS['navy_900'])
    bg.line.fill.background()
    
    # Main message
    add_text_box(
        slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.5),
        slide_data.get('title', '让数据真正从"专家工具"\n走向"人人可用"'),
        font_size=Pt(36), font_color=COLORS['white'], bold=True
    )
    
    # CTA or contact
    cta = slide_data.get('cta', '') or slide_data.get('contact_name', '')
    if cta:
        add_text_box(
            slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.8),
            cta,
            font_size=Pt(20), font_color=COLORS['slate_400']
        )
    
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.5), Inches(3), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(COLORS['blue'])
    line.line.fill.background()
    
    # Company name
    contact = slide_data.get('contact_name', '数势科技 · DigitForce')
    add_text_box(
        slide, Inches(0.8), Inches(6), Inches(5), Inches(0.5),
        contact,
        font_size=Pt(16), font_color=COLORS['slate_400']
    )
    
    # Website
    website = slide_data.get('contact_email', '')
    if website:
        add_text_box(
            slide, Inches(6), Inches(6), Inches(5), Inches(0.5),
            website,
            font_size=Pt(14), font_color=COLORS['slate_400']
        )
    
    # Page number
    add_text_box(
        slide, Inches(11.5), Inches(0.5), Inches(1.5), Inches(0.8),
        f"{slide_num:02d}",
        font_size=Pt(36), font_color=COLORS['navy_700'], bold=True, align='right'
    )
    
    return slide


def create_presentation(slides, output_path="output.pptx"):
    """Create professional PPTX presentation"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for idx, slide_data in enumerate(slides):
        slide_num = idx + 1
        template = detect_template(slide_data)
        
        if template == 'COVER':
            create_cover_slide(prs, slide_data, slide_num)
        elif template == 'PROBLEM':
            create_problem_slide(prs, slide_data, slide_num)
        elif template == 'SOLUTION':
            create_solution_slide(prs, slide_data, slide_num)
        elif template == 'MATRIX':
            create_matrix_slide(prs, slide_data, slide_num)
        elif template == 'CLOSING':
            create_closing_slide(prs, slide_data, slide_num)
        else:
            create_content_slide(prs, slide_data, slide_num)
    
    prs.save(output_path)
    
    return {
        "success": True,
        "output": output_path,
        "total_slides": len(prs.slides),
        "message": f"Created {len(prs.slides)} professional slides"
    }


def main():
    parser = argparse.ArgumentParser(description='Generate professional PPTX from content')
    parser.add_argument('--output', '-o', required=True, help='Output PPTX file path')
    parser.add_argument('--content', help='JSON file with slides array')
    parser.add_argument('--slides', help='JSON string with slides array')
    parser.add_argument('--title', default='数势科技', help='Presentation title')
    
    args = parser.parse_args()
    
    # Get slides
    slides = []
    
    if args.slides:
        if os.path.exists(args.slides):
            with open(args.slides, 'r', encoding='utf-8') as f:
                slides = json.load(f)
        else:
            try:
                slides = json.loads(args.slides)
            except:
                slides = [{"title": "Content", "content": args.slides}]
    elif args.content:
        with open(args.content, 'r', encoding='utf-8') as f:
            data = json.load(f)
            if isinstance(data, list):
                slides = data
            else:
                slides = data.get('slides', [data])
    else:
        try:
            data = json.load(sys.stdin)
            if isinstance(data, list):
                slides = data
            else:
                slides = data.get('slides', [data])
        except:
            print(json.dumps({"success": False, "error": "No content provided"}))
            sys.exit(1)
    
    result = create_presentation(slides=slides, output_path=args.output)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
