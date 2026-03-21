#!/usr/bin/env python3
# Part of doc2slides skill.

#!/usr/bin/env python3
"""
Validate generated PPTX quality.
Checks content completeness, formatting, readability, visual consistency.

Usage:
  python validate.py --input presentation.pptx --output report.json
"""

import sys
import json
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def check_content_completeness(prs) -> list:
    """Check if each slide has meaningful content."""
    issues = []
    
    for i, slide in enumerate(prs.slides, 1):
        text_content = []
        shape_count = 0
        
        for shape in slide.shapes:
            shape_count += 1
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    text_content.append(text)
        
        # Check for empty slides
        if shape_count == 0:
            issues.append({
                'slide': i,
                'type': 'empty_slide',
                'message': f'Slide {i}: 空白幻灯片，无任何元素'
            })
        elif not text_content:
            issues.append({
                'slide': i,
                'type': 'no_text',
                'message': f'Slide {i}: 无文本内容'
            })
        
        # Check for very short content
        total_text = ' '.join(text_content)
        if len(total_text) < 20:
            issues.append({
                'slide': i,
                'type': 'too_short',
                'message': f'Slide {i}: 内容过短 ({len(total_text)} 字符)'
            })
    
    return issues


def check_text_readability(prs) -> list:
    """Check text readability issues."""
    issues = []
    
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    
                    # Check text length per slide
                    if len(text) > 200:
                        issues.append({
                            'slide': i,
                            'type': 'too_long',
                            'message': f'Slide {i}: 单段文字过长 ({len(text)} 字符)，建议拆分'
                        })
                    
                    # Check font size
                    if para.font.size:
                        size_pt = para.font.size.pt
                        if size_pt < 12:
                            issues.append({
                                'slide': i,
                                'type': 'font_too_small',
                                'message': f'Slide {i}: 字体过小 ({size_pt}pt)，建议至少 12pt'
                            })
    
    return issues


def check_visual_consistency(prs) -> list:
    """Check visual consistency across slides."""
    issues = []
    
    # Collect colors and fonts
    colors = set()
    fonts = set()
    font_sizes = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.name:
                        fonts.add(para.font.name)
                    if para.font.size:
                        font_sizes.append(para.font.size.pt)
                    
                    for run in para.runs:
                        if run.font.color.rgb:
                            colors.add(str(run.font.color.rgb))
    
    # Check for too many fonts
    if len(fonts) > 3:
        issues.append({
            'slide': 0,
            'type': 'too_many_fonts',
            'message': f'使用了 {len(fonts)} 种字体，建议控制在 2-3 种以内'
        })
    
    # Check for inconsistent font sizes
    if font_sizes:
        size_range = max(font_sizes) - min(font_sizes)
        if size_range > 24:
            issues.append({
                'slide': 0,
                'type': 'inconsistent_sizes',
                'message': f'字体大小差异过大 ({min(font_sizes)}-{max(font_sizes)}pt)，建议统一层级'
            })
    
    return issues


def check_layout_alignment(prs) -> list:
    """Check layout alignment issues."""
    issues = []
    
    for i, slide in enumerate(prs.slides, 1):
        # Check for overlapping shapes
        shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                shapes.append({
                    'x': shape.left / 914400,
                    'y': shape.top / 914400,
                    'w': shape.width / 914400,
                    'h': shape.height / 914400,
                    'text': shape.text_frame.text[:30]
                })
        
        # Check overlap
        for j, s1 in enumerate(shapes):
            for s2 in shapes[j+1:]:
                # Check if shapes overlap
                if (s1['x'] < s2['x'] + s2['w'] and s1['x'] + s1['w'] > s2['x'] and
                    s1['y'] < s2['y'] + s2['h'] and s1['y'] + s1['h'] > s2['y']):
                    issues.append({
                        'slide': i,
                        'type': 'overlap',
                        'message': f'Slide {i}: 元素重叠 ({s1["text"][:10]}... / {s2["text"][:10]}...)'
                    })
    
    return issues


def validate_pptx(pptx_path: str) -> dict:
    """
    Validate PPTX file and return report.
    """
    if not HAS_PPTX:
        return {
            'passed': False,
            'error': 'python-pptx not installed'
        }
    
    try:
        prs = Presentation(pptx_path)
        
        all_issues = []
        all_issues.extend(check_content_completeness(prs))
        all_issues.extend(check_text_readability(prs))
        all_issues.extend(check_visual_consistency(prs))
        all_issues.extend(check_layout_alignment(prs))
        
        # Calculate quality score
        total_checks = len(prs.slides) * 4  # 4 checks per slide
        failed_checks = len(all_issues)
        score = max(0, 100 - (failed_checks * 5))
        
        return {
            'passed': len(all_issues) == 0,
            'score': score,
            'slide_count': len(prs.slides),
            'issues_count': len(all_issues),
            'issues': all_issues,
            'summary': {
                'empty_slides': len([i for i in all_issues if i['type'] == 'empty_slide']),
                'content_issues': len([i for i in all_issues if i['type'] in ['no_text', 'too_short', 'too_long']]),
                'formatting_issues': len([i for i in all_issues if i['type'] in ['font_too_small', 'too_many_fonts']]),
                'layout_issues': len([i for i in all_issues if i['type'] == 'overlap'])
            }
        }
    except Exception as e:
        return {
            'passed': False,
            'error': str(e)
        }


def main():
    parser = argparse.ArgumentParser(description='Validate PPTX quality')
    parser.add_argument('--input', required=True, help='Input PPTX file')
    parser.add_argument('--output', help='Output JSON report')
    
    args = parser.parse_args()
    
    report = validate_pptx(args.input)
    
    output_json = json.dumps(report, ensure_ascii=False, indent=2)
    
    if args.output:
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(output_json)
        print(f"✓ Validation report saved to {args.output}")
    else:
        print(output_json)
    
    # Print summary
    if report.get('passed'):
        print(f"\n✅ PPT 质量验证通过！")
    else:
        print(f"\n⚠️ 发现 {report.get('issues_count', 0)} 个问题：")
        for issue in report.get('issues', [])[:5]:
            print(f"  - {issue['message']}")


if __name__ == "__main__":
    main()
