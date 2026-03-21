#!/usr/bin/env python3
# Part of doc2slides skill.

#!/usr/bin/env python3
"""提取 PPT 结构"""

from pptx import Presentation
import sys

def extract_ppt_structure(ppt_path):
    prs = Presentation(ppt_path)
    
    print(f"PPT 文件: {ppt_path}")
    print(f"总页数: {len(prs.slides)}")
    print(f"幻灯片尺寸: {prs.slide_width.inches} x {prs.slide_height.inches} 英寸")
    print()
    
    for i, slide in enumerate(prs.slides, 1):
        print(f"--- 幻灯片 {i} ---")
        
        # 提取所有文本
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        
        if texts:
            for text in texts[:5]:  # 只显示前5个文本框
                print(f"  {text[:100]}")
        else:
            print("  (无文本，可能是图片型幻灯片)")
        
        print()

if __name__ == "__main__":
    ppt_path = sys.argv[1] if len(sys.argv) > 1 else ""
    if not ppt_path:
        print("Usage: python extract_ppt_structure.py <pptx_path>")
        sys.exit(1)
    extract_ppt_structure(ppt_path)
