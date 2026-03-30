#!/usr/bin/env python3
import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path


def run(cmd):
    return subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)


def available_tess_langs():
    r = run(["tesseract", "--list-langs"])
    if r.returncode != 0:
        return []
    lines = [x.strip() for x in r.stdout.splitlines() if x.strip()]
    if lines and lines[0].lower().startswith("list of available"):
        lines = lines[1:]
    return lines


def extract_pdf(path: Path, page_limit: int):
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    page_count = len(reader.pages)
    texts = []
    for i, page in enumerate(reader.pages[:page_limit], start=1):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        txt = txt.strip()
        if txt:
            texts.append(f"\n## Page {i}\n{txt}\n")
    joined = "\n".join(texts).strip()
    used_ocr = False
    if len(joined) >= 200:
        return {
            "type": "pdf",
            "page_count": page_count,
            "extraction": "text",
            "text": joined,
        }

    langs = available_tess_langs()
    tess_lang = None
    if "chi_sim" in langs and "eng" in langs:
        tess_lang = "chi_sim+eng"
    elif "chi_sim" in langs:
        tess_lang = "chi_sim"
    elif "eng" in langs:
        tess_lang = "eng"

    if not shutil.which("pdftoppm") or not shutil.which("tesseract") or not tess_lang:
        return {
            "type": "pdf",
            "page_count": page_count,
            "extraction": "text-empty",
            "text": joined,
            "warning": "No usable OCR pipeline available; install poppler-utils, tesseract-ocr, and tesseract language packs (prefer chi_sim+eng).",
        }

    used_ocr = True
    with tempfile.TemporaryDirectory(prefix="office-doc-pdf-") as td:
        prefix = os.path.join(td, "page")
        cmd = ["pdftoppm", "-f", "1", "-l", str(min(page_limit, page_count)), "-r", "180", "-png", str(path), prefix]
        r = run(cmd)
        if r.returncode != 0:
            return {
                "type": "pdf",
                "page_count": page_count,
                "extraction": "ocr-failed",
                "text": joined,
                "warning": r.stderr.strip()[:1000],
            }
        page_texts = []
        for img in sorted(Path(td).glob("page-*.png")):
            rr = run(["tesseract", str(img), "stdout", "-l", tess_lang])
            txt = (rr.stdout or "").strip()
            if txt:
                page_texts.append(f"\n## {img.stem}\n{txt}\n")
        ocr_text = "\n".join(page_texts).strip()
        combined = ocr_text if len(ocr_text) > len(joined) else joined
        return {
            "type": "pdf",
            "page_count": page_count,
            "extraction": "ocr" if used_ocr else "text",
            "ocr_lang": tess_lang,
            "text": combined,
        }


def extract_docx(path: Path):
    from docx import Document

    doc = Document(str(path))
    parts = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)
    for table_idx, table in enumerate(doc.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"\n[Table {table_idx}]\n" + "\n".join(rows))
    return {"type": "docx", "text": "\n\n".join(parts).strip()}


def extract_doc(path: Path):
    for cmd in (["antiword", str(path)], ["catdoc", str(path)]):
        if shutil.which(cmd[0]):
            r = run(cmd)
            if r.returncode == 0 and r.stdout.strip():
                return {"type": "doc", "extraction": cmd[0], "text": r.stdout.strip()}
    if shutil.which("libreoffice"):
        with tempfile.TemporaryDirectory(prefix="office-doc-doc-") as td:
            r = run(["libreoffice", "--headless", "--convert-to", "docx", "--outdir", td, str(path)])
            out = Path(td) / (path.stem + ".docx")
            if out.exists():
                data = extract_docx(out)
                data["type"] = "doc"
                data["extraction"] = "libreoffice->docx"
                return data
            return {"type": "doc", "warning": (r.stderr or r.stdout).strip()[:1000], "text": ""}
    return {"type": "doc", "warning": "No DOC extractor found (antiword/catdoc/libreoffice).", "text": ""}


def sheet_to_lines(ws, row_limit):
    rows = []
    max_row = min(ws.max_row or 0, row_limit)
    max_col = ws.max_column or 0
    for r in ws.iter_rows(min_row=1, max_row=max_row, min_col=1, max_col=max_col, values_only=True):
        vals = ["" if v is None else str(v).strip() for v in r]
        if any(vals):
            rows.append(" | ".join(vals))
    return rows


def extract_xlsx(path: Path, row_limit: int):
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        lines = sheet_to_lines(ws, row_limit)
        parts.append(f"\n# Sheet: {name}\n")
        if lines:
            parts.extend(lines)
        else:
            parts.append("(empty)")
    return {"type": "xlsx", "sheets": wb.sheetnames, "text": "\n".join(parts).strip()}


def extract_xls(path: Path, row_limit: int):
    if shutil.which("libreoffice"):
        with tempfile.TemporaryDirectory(prefix="office-doc-xls-") as td:
            r = run(["libreoffice", "--headless", "--convert-to", "xlsx", "--outdir", td, str(path)])
            out = Path(td) / (path.stem + ".xlsx")
            if out.exists():
                data = extract_xlsx(out, row_limit)
                data["type"] = "xls"
                data["extraction"] = "libreoffice->xlsx"
                return data
            return {"type": "xls", "warning": (r.stderr or r.stdout).strip()[:1000], "text": ""}
    return {"type": "xls", "warning": "No XLS converter found (libreoffice).", "text": ""}


def extract_pptx(path: Path):
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            txt = getattr(shape, "text", "")
            txt = (txt or "").strip()
            if txt:
                lines.append(txt)
        if slide.has_notes_slide:
            notes_texts = []
            for shape in slide.notes_slide.shapes:
                txt = getattr(shape, "text", "")
                txt = (txt or "").strip()
                if txt:
                    notes_texts.append(txt)
            if notes_texts:
                lines.append("[Notes]\n" + "\n".join(notes_texts))
        if lines:
            parts.append(f"\n# Slide {idx}\n" + "\n".join(lines))
        else:
            parts.append(f"\n# Slide {idx}\n(no extractable text)")
    return {"type": "pptx", "slide_count": len(prs.slides), "text": "\n\n".join(parts).strip()}


def extract_ppt(path: Path):
    if shutil.which("libreoffice"):
        with tempfile.TemporaryDirectory(prefix="office-doc-ppt-") as td:
            r = run(["libreoffice", "--headless", "--convert-to", "pptx", "--outdir", td, str(path)])
            out = Path(td) / (path.stem + ".pptx")
            if out.exists():
                data = extract_pptx(out)
                data["type"] = "ppt"
                data["extraction"] = "libreoffice->pptx"
                return data
            return {"type": "ppt", "warning": (r.stderr or r.stdout).strip()[:1000], "text": ""}
    return {"type": "ppt", "warning": "No PPT converter found (libreoffice).", "text": ""}


def main():
    ap = argparse.ArgumentParser(description="Extract text from office documents for OpenClaw skills.")
    ap.add_argument("file")
    ap.add_argument("--page-limit", type=int, default=20)
    ap.add_argument("--row-limit", type=int, default=50)
    ap.add_argument("--max-chars", type=int, default=50000)
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    path = Path(args.file).expanduser().resolve()
    if not path.exists():
        print(json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False))
        sys.exit(2)

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        data = extract_pdf(path, args.page_limit)
    elif suffix == ".docx":
        data = extract_docx(path)
    elif suffix == ".doc":
        data = extract_doc(path)
    elif suffix == ".xlsx":
        data = extract_xlsx(path, args.row_limit)
    elif suffix == ".xls":
        data = extract_xls(path, args.row_limit)
    elif suffix == ".pptx":
        data = extract_pptx(path)
    elif suffix == ".ppt":
        data = extract_ppt(path)
    else:
        data = {"error": f"Unsupported extension: {suffix}"}

    data["file"] = str(path)
    text = data.get("text", "") or ""
    if len(text) > args.max_chars:
        data["text"] = text[: args.max_chars] + "\n\n[TRUNCATED]"
        data["truncated"] = True
    if args.json:
        print(json.dumps(data, ensure_ascii=False, indent=2))
    else:
        meta = {k: v for k, v in data.items() if k != "text"}
        print(json.dumps(meta, ensure_ascii=False, indent=2))
        print("\n---TEXT---\n")
        print(data.get("text", ""))


if __name__ == "__main__":
    main()
