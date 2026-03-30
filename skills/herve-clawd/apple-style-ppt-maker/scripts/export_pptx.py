#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "python-pptx>=1.0.2",
#   "pillow>=10.0.0",
# ]
# ///
"""Export generated slide images to a 16:9 PPTX file."""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))


def natural_key(path: Path) -> int:
    match = re.search(r"(\d+)", path.stem)
    return int(match.group(1)) if match else 0


def import_pptx_modules():
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise RuntimeError("python-pptx is required. Run: pip install python-pptx pillow") from exc
    return Presentation, Inches


def enable_webp_support() -> None:
    """
    Patch python-pptx runtime mappings to accept WebP images.

    python-pptx 1.0.x may reject WebP in add_picture() without this patch.
    """
    try:
        from pptx.opc import spec as opc_spec
        from pptx.parts.image import Image as PptxImage
        from pptx.util import lazyproperty
    except Exception:
        return

    opc_spec.image_content_types.setdefault("webp", "image/webp")

    @lazyproperty  # type: ignore[misc]
    def ext(self) -> str:
        mapping = {
            "BMP": "bmp",
            "GIF": "gif",
            "JPEG": "jpg",
            "PNG": "png",
            "TIFF": "tiff",
            "WEBP": "webp",
            "WMF": "wmf",
        }
        if self._format not in mapping:
            expected = ", ".join(mapping.keys())
            raise ValueError(f"unsupported image format, expected one of: {expected}, got '{self._format}'")
        return mapping[self._format]

    PptxImage.ext = ext  # type: ignore[assignment]


def collect_slide_images(images_dir: Path, glob_pattern: str | None) -> list[Path]:
    if glob_pattern:
        return sorted([p for p in images_dir.glob(glob_pattern) if p.is_file()], key=natural_key)

    for ext in ("webp", "png", "jpg", "jpeg"):
        files = sorted([p for p in images_dir.glob(f"slide-*.{ext}") if p.is_file()], key=natural_key)
        if files:
            return files

    return sorted([p for p in images_dir.glob("slide-*.*") if p.is_file()], key=natural_key)


def export_pptx(images_dir: Path, output_pptx: Path, glob_pattern: str | None = None) -> None:
    Presentation, Inches = import_pptx_modules()
    enable_webp_support()
    images = collect_slide_images(images_dir, glob_pattern=glob_pattern)
    if not images:
        raise SystemExit(f"No slide images found in: {images_dir}")

    deck = Presentation()
    deck.slide_width = Inches(13.333)
    deck.slide_height = Inches(7.5)
    blank = deck.slide_layouts[6]

    for image in images:
        slide = deck.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), left=Inches(0), top=Inches(0), width=deck.slide_width, height=deck.slide_height)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    deck.save(str(output_pptx))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Export slide images to PPTX")
    parser.add_argument("--images-dir", required=True, help="Directory containing slide images")
    parser.add_argument("--output-pptx", required=True, help="Target PPTX file")
    parser.add_argument("--glob", default=None, help="Optional glob pattern, e.g. 'slide-*.webp'")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    images_dir = Path(args.images_dir).resolve()
    output_pptx = Path(args.output_pptx).resolve()

    if not images_dir.exists():
        print(f"[ERROR] Images directory not found: {images_dir}")
        return 2

    try:
        export_pptx(images_dir, output_pptx, glob_pattern=args.glob)
    except Exception as exc:
        print(f"[ERROR] Export failed: {exc}")
        return 1

    print(f"[DONE] PPTX exported: {output_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
