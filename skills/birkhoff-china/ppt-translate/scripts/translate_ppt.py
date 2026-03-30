#!/usr/bin/env python3
"""
PPT Translation Tool: Chinese to English
Translates PPTX presentations while preserving all non-text content.

Usage:
    python translate_ppt.py <input.pptx> [output.pptx]
    python translate_ppt.py <input.pptx> --font "Arial" --model qwen2.5:14b --verbose
"""

import argparse
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from time import sleep
from typing import List, Optional, Tuple, Dict, Any

# Third-party imports with graceful fallback
try:
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.dml.color import RGBColor
except ImportError:
    print("Error: python-pptx is required. Install with: pip install python-pptx")
    sys.exit(1)

try:
    import requests
except ImportError:
    print("Error: requests is required. Install with: pip install requests")
    sys.exit(1)


# =============================================================================
# Configuration Constants
# =============================================================================

CJK_FONTS = {
    'SimSun', '宋体',
    'SimHei', '黑体',
    'Microsoft YaHei', '微软雅黑',
    'KaiTi', '楷体',
    'FangSong', '仿宋',
    'STSong', 'STHeiti', 'STKaiti', 'STFangsong',
    'NSimSun', '新宋体',
    'PMingLiU', 'MingLiU',
    'DengXian', '等线',
    'Source Han Sans', 'Noto Sans CJK', 'Source Han Serif',
    'WenQuanYi Micro Hei', 'WenQuanYi Zen Hei',
}

# Unicode ranges for CJK characters
CJK_RANGES = [
    (0x4E00, 0x9FFF),   # CJK Unified Ideographs
    (0x3400, 0x4DBF),   # CJK Extension A
    (0xF900, 0xFAFF),   # CJK Compatibility Ideographs
    (0x20000, 0x2A6DF), # CJK Extension B
    (0x2A700, 0x2B73F), # CJK Extension C
    (0x2B740, 0x2B81F), # CJK Extension D
    (0x3000, 0x303F),   # CJK Symbols and Punctuation
    (0xFF00, 0xFFEF),   # Halfwidth and Fullwidth Forms
]

DEFAULT_FONT = "Calibri"
DEFAULT_MODEL = "qwen2.5:14b"
DEFAULT_API_BASE = "http://localhost:11434/v1"
DEFAULT_BATCH_SIZE = 20
MAX_RETRIES = 3

SYSTEM_PROMPT = """You are a professional translator specializing in business and technical content. 
Translate Chinese text to English following these rules:
1. Maintain original meaning, tone, and intent
2. Use professional business English for formal content
3. Preserve formatting markers (bullet points, numbering) in the translation
4. Do not translate proper nouns unless they have standard English equivalents
5. For mixed Chinese/English text, only translate the Chinese portions
6. Return ONLY the translations, one per line, corresponding to each input line
7. Do not add numbering, explanations, or any other text

Input format: Numbered list of text segments
Output format: Same numbered list with English translations only"""


# =============================================================================
# Data Classes
# =============================================================================

@dataclass
class TextSegment:
    """Represents a single text run to be translated."""
    text: str
    slide_idx: int
    shape_path: str
    para_idx: int
    run_idx: int
    is_bold: bool = False
    is_italic: bool = False
    font_name: Optional[str] = None
    font_size: Optional[Any] = None  # Pt object
    translated: str = ""
    
    def __post_init__(self):
        # Store original text length for overflow estimation
        self.original_length = len(self.text)


# =============================================================================
# Utility Functions
# =============================================================================

def contains_chinese(text: str) -> bool:
    """Check if text contains any Chinese/CJK characters."""
    if not text or not text.strip():
        return False
    for char in text:
        code = ord(char)
        for start, end in CJK_RANGES:
            if start <= code <= end:
                return True
    return False


def is_cjk_font(font_name: Optional[str]) -> bool:
    """Check if font name is a CJK font that should be replaced."""
    if not font_name:
        return True  # Default to replacing unset fonts
    font_lower = font_name.lower()
    if font_name in CJK_FONTS:
        return True
    # Check for common CJK font patterns
    cjk_patterns = ['song', 'hei', 'kai', 'fang', 'ming', 'gothic', 'yuan', 'meiryo', 'yugothic']
    return any(pattern in font_lower for pattern in cjk_patterns)


def parse_numbered_response(response_text: str, expected_count: int) -> List[str]:
    """Parse numbered response from LLM into list of translations."""
    lines = response_text.strip().split('\n')
    translations = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Remove numbering patterns like "1. ", "1) ", "1 - "
        cleaned = re.sub(r'^\d+[\.\)\-\s]+', '', line)
        translations.append(cleaned)
    
    # Handle case where response has fewer lines than expected
    while len(translations) < expected_count:
        translations.append("")
    
    return translations[:expected_count]


def get_output_path(input_path: str, output_path: Optional[str] = None) -> str:
    """Generate output path if not provided."""
    if output_path:
        return output_path
    
    path = Path(input_path)
    stem = path.stem
    suffix = path.suffix
    parent = path.parent
    
    # Remove _cn or _zh suffix if present, then add _en
    if stem.endswith(('_cn', '_zh', '_CN', '_ZH')):
        stem = stem[:-3]
    
    return str(parent / f"{stem}_en{suffix}")


# =============================================================================
# LLM API Helpers
# =============================================================================

def check_api_connectivity(api_base: str) -> bool:
    """Check if the LLM API endpoint is reachable."""
    try:
        # Try a simple GET to the base URL
        resp = requests.get(api_base.rstrip('/v1').rstrip('/'), timeout=5)
        return resp.status_code < 500
    except requests.ConnectionError:
        return False
    except Exception:
        return False


# =============================================================================
# Text Extraction
# =============================================================================

def extract_from_text_frame(
    text_frame,
    slide_idx: int,
    shape_path: str,
    segments: List[TextSegment]
) -> None:
    """Extract all text runs from a text frame."""
    for para_idx, para in enumerate(text_frame.paragraphs):
        for run_idx, run in enumerate(para.runs):
            text = run.text
            if text and text.strip() and contains_chinese(text):
                segment = TextSegment(
                    text=text,
                    slide_idx=slide_idx,
                    shape_path=shape_path,
                    para_idx=para_idx,
                    run_idx=run_idx,
                    is_bold=bool(run.font.bold),
                    is_italic=bool(run.font.italic),
                    font_name=run.font.name,
                    font_size=run.font.size
                )
                segments.append(segment)


def extract_from_shape(
    shape,
    slide_idx: int,
    shape_path: str,
    segments: List[TextSegment]
) -> None:
    """Recursively extract text from a shape and its children."""
    # Handle regular shapes with text frames
    if shape.has_text_frame:
        extract_from_text_frame(shape.text_frame, slide_idx, shape_path, segments)
    
    # Handle grouped shapes recursively
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_idx, child in enumerate(shape.shapes):
            child_path = f"{shape_path}/group[{child_idx}]"
            extract_from_shape(child, slide_idx, child_path, segments)
    
    # Handle tables
    if shape.has_table:
        table = shape.table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_path = f"{shape_path}/table[{row_idx},{col_idx}]"
                if cell.text_frame:
                    extract_from_text_frame(cell.text_frame, slide_idx, cell_path, segments)


def extract_from_notes(
    slide,
    slide_idx: int,
    segments: List[TextSegment]
) -> None:
    """Extract text from slide notes."""
    if slide.has_notes_slide and slide.notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            extract_from_text_frame(
                notes_frame,
                slide_idx,
                "notes",
                segments
            )


def extract_all_texts(prs: Presentation) -> List[TextSegment]:
    """Extract all Chinese text segments from the presentation."""
    segments = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            shape_path = f"slide[{slide_idx}]/shape[{shape_idx}]"
            extract_from_shape(shape, slide_idx, shape_path, segments)
        
        # Extract notes
        extract_from_notes(slide, slide_idx, segments)
    
    return segments


# =============================================================================
# Translation
# =============================================================================

def translate_batch(
    segments: List[TextSegment],
    model: str,
    api_base: str,
    api_key: Optional[str] = None,
    verbose: bool = False
) -> List[str]:
    """Translate a batch of text segments using OpenAI-compatible API."""
    if not segments:
        return []

    texts = [seg.text for seg in segments]
    numbered_texts = "\n".join(f"{i+1}. {t}" for i, t in enumerate(texts))

    url = f"{api_base}/chat/completions"

    payload = {
        "model": model,
        "messages": [
            {
                "role": "system",
                "content": SYSTEM_PROMPT
            },
            {
                "role": "user",
                "content": f"Translate the following {len(texts)} Chinese text segments to English. Return ONLY the translations, one per line, with matching line numbers:\n\n{numbered_texts}"
            }
        ],
        "temperature": 0.3,
        "stream": False
    }

    headers = {"Content-Type": "application/json"}
    if api_key:
        headers["Authorization"] = f"Bearer {api_key}"

    if verbose:
        print(f"    Sending batch of {len(segments)} segments to LLM...")

    for attempt in range(MAX_RETRIES):
        try:
            response = requests.post(url, json=payload, headers=headers, timeout=120)
            response.raise_for_status()
            result = response.json()
            translated_text = result["choices"][0]["message"]["content"]
            translations = parse_numbered_response(translated_text, len(segments))

            if verbose:
                print(f"    Received {len(translations)} translations")

            return translations

        except requests.ConnectionError as e:
            if attempt < MAX_RETRIES - 1:
                wait_time = 2 ** attempt
                print(f"    Connection error, waiting {wait_time}s before retry...")
                sleep(wait_time)
            else:
                print(f"    Warning: Cannot connect to LLM API after retries. Check your endpoint.")
                return [seg.text for seg in segments]  # Return original texts

        except requests.Timeout as e:
            if attempt < MAX_RETRIES - 1:
                wait_time = 2 ** attempt
                print(f"    Request timed out, waiting {wait_time}s before retry...")
                sleep(wait_time)
            else:
                print(f"    Warning: LLM request timed out. The model may be loading. Try again.")
                return [seg.text for seg in segments]

        except requests.HTTPError as e:
            error_msg = str(e)
            if "404" in error_msg:
                print(f"    Error: Model '{model}' not found or endpoint not available.")
                return [seg.text for seg in segments]
            elif attempt < MAX_RETRIES - 1:
                print(f"    HTTP error: {e}, retrying...")
                sleep(1)
            else:
                print(f"    Warning: HTTP error after retries: {e}")
                return [seg.text for seg in segments]

        except Exception as e:
            print(f"    Warning: Unexpected error during translation: {e}")
            return [seg.text for seg in segments]

    return [seg.text for seg in segments]


def translate_segments(
    segments: List[TextSegment],
    model: str,
    api_base: str,
    api_key: Optional[str],
    batch_size: int,
    verbose: bool = False
) -> None:
    """Translate all segments in batches."""
    if not segments:
        return

    total = len(segments)
    print(f"Translating {total} text segments in batches of {batch_size}...")

    for i in range(0, total, batch_size):
        batch = segments[i:i + batch_size]
        batch_num = i // batch_size + 1
        total_batches = (total + batch_size - 1) // batch_size

        print(f"  Batch {batch_num}/{total_batches} ({len(batch)} segments)...")
        translations = translate_batch(batch, model, api_base, api_key, verbose)

        # Assign translations back to segments
        for segment, translation in zip(batch, translations):
            segment.translated = translation


# =============================================================================
# Text Application
# =============================================================================

def locate_shape(slide, shape_path: str):
    """Locate a shape by its path string."""
    # Parse the path to navigate to the shape
    # Format: slide[X]/shape[Y]/group[Z]/table[R,C] or notes
    
    if shape_path == "notes":
        if slide.has_notes_slide and slide.notes_slide:
            return slide.notes_slide.notes_text_frame
        return None
    
    parts = shape_path.split('/')
    current = None
    
    for part in parts:
        if part.startswith("slide["):
            continue  # Skip slide part, we already have the slide
        
        elif part.startswith("shape["):
            idx = int(part[6:-1])  # Extract index from "shape[N]"
            if idx < len(slide.shapes):
                current = slide.shapes[idx]
        
        elif part.startswith("group["):
            idx = int(part[6:-1])
            if current and current.shape_type == MSO_SHAPE_TYPE.GROUP:
                shapes_list = list(current.shapes)
                if idx < len(shapes_list):
                    current = shapes_list[idx]
        
        elif part.startswith("table["):
            # Extract row,col from "table[R,C]"
            coords = part[6:-1].split(',')
            row_idx = int(coords[0])
            col_idx = int(coords[1])
            if current and current.has_table:
                table = current.table
                if row_idx < len(table.rows):
                    row = table.rows[row_idx]
                    if col_idx < len(row.cells):
                        current = row.cells[col_idx].text_frame
    
    return current


def apply_translation(
    prs: Presentation,
    segment: TextSegment,
    target_font: str
) -> None:
    """Apply a single translation to the presentation."""
    slide = prs.slides[segment.slide_idx]
    target = locate_shape(slide, segment.shape_path)
    
    if target is None:
        return
    
    # Handle notes text frame directly
    if segment.shape_path == "notes":
        text_frame = target
    elif hasattr(target, 'text_frame'):
        text_frame = target.text_frame
    else:
        text_frame = target
    
    # Access the specific paragraph and run
    if segment.para_idx < len(text_frame.paragraphs):
        para = text_frame.paragraphs[segment.para_idx]
        if segment.run_idx < len(para.runs):
            run = para.runs[segment.run_idx]
            
            # Replace text
            run.text = segment.translated
            
            # Check if parent shape is a title placeholder
            is_title = is_title_shape(target) if hasattr(target, 'placeholder_format') else False
            
            # Replace font if it's a CJK font
            if is_cjk_font(run.font.name):
                run.font.name = target_font
                # If original was bold (likely a title) OR shape is a title placeholder, keep/make it bold
                if segment.is_bold or is_title:
                    run.font.bold = True
            
            # Apply title formatting for title placeholders (force bold Calibri)
            if is_title:
                run.font.name = target_font
                run.font.bold = True
            
            # Adjust text box size after translation (only for non-table, non-group shapes)
            if not segment.shape_path.startswith("notes") and "table[" not in segment.shape_path:
                adjust_text_box_size(target, segment.text, segment.translated, is_title)


def is_title_shape(shape) -> bool:
    """Check if a shape is a title placeholder."""
    ph = getattr(shape, "placeholder_format", None)
    if ph is None:
        return False
    try:
        return ph.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE)
    except Exception:
        return False


def adjust_text_box_size(shape, original_text: str, translated_text: str, is_title: bool = False) -> None:
    """Adjust text box sizing for English text."""
    if not hasattr(shape, 'text_frame'):
        return
    
    text_frame = shape.text_frame
    
    # Enable word wrap
    text_frame.word_wrap = True
    
    # Check if text is significantly longer
    if len(translated_text) > len(original_text) * 1.5:
        # Try to enable auto-size if available
        try:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except:
            pass
        
        # Attempt to expand shape width by up to 20% for freestanding text boxes
        # Skip if shape is grouped or in a table
        if hasattr(shape, 'width') and hasattr(shape, 'parent'):
            # Check if shape is not part of a group or table
            parent = getattr(shape, 'parent', None)
            is_grouped = False
            is_in_table = False
            
            # Check if parent is a group shape
            if parent is not None:
                parent_type = getattr(parent, 'shape_type', None)
                if parent_type == MSO_SHAPE_TYPE.GROUP:
                    is_grouped = True
            
            # Check if we're in a table cell by looking at the shape path
            # (This is handled by not calling adjust_text_box_size for table cells)
            
            if not is_grouped and not is_in_table:
                current_width = shape.width
                # Max slide width is ~9144000 EMU (10 inches)
                max_width = 9144000
                new_width = min(int(current_width * 1.2), max_width)
                if new_width > current_width:
                    try:
                        shape.width = new_width
                    except:
                        pass
        
        # Optionally reduce font size as last resort fallback
        # Only for non-title text and if font size is reasonably large
        if not is_title:
            for para in text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size > Pt(10):
                        try:
                            new_size = max(Pt(10), run.font.size - Pt(1))
                            run.font.size = new_size
                        except:
                            pass


def apply_translations(
    prs: Presentation,
    segments: List[TextSegment],
    target_font: str,
    verbose: bool = False
) -> int:
    """Apply all translations to the presentation."""
    count = 0
    
    for segment in segments:
        if segment.translated:
            apply_translation(prs, segment, target_font)
            count += 1
    
    return count


# =============================================================================
# Font Adjustment Pass
# =============================================================================

def adjust_fonts_in_text_frame(text_frame, target_font: str, is_title: bool = False) -> None:
    """Adjust fonts in a text frame to target font."""
    for para in text_frame.paragraphs:
        for run in para.runs:
            if is_cjk_font(run.font.name):
                run.font.name = target_font
                # Force bold for title placeholders
                if is_title:
                    run.font.bold = True


def adjust_fonts_in_shape(shape, target_font: str) -> None:
    """Recursively adjust fonts in a shape and its children."""
    is_title = is_title_shape(shape)
    
    if shape.has_text_frame:
        adjust_fonts_in_text_frame(shape.text_frame, target_font, is_title)
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            adjust_fonts_in_shape(child, target_font)
    
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    adjust_fonts_in_text_frame(cell.text_frame, target_font, is_title)


def adjust_all_fonts(prs: Presentation, target_font: str) -> None:
    """Adjust all CJK fonts to target font throughout the presentation."""
    for slide in prs.slides:
        for shape in slide.shapes:
            adjust_fonts_in_shape(shape, target_font)
        
        # Adjust notes fonts
        if slide.has_notes_slide and slide.notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                adjust_fonts_in_text_frame(notes_frame, target_font)


# =============================================================================
# Main Function
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Translate Chinese PowerPoint presentations to English",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Use with local Ollama (default)
  python translate_ppt.py input.pptx

  # Use with Qoderwork or other OpenAI-compatible endpoint
  python translate_ppt.py input.pptx --api-base http://localhost:PORT/v1 --model model-name

  # Use with cloud API
  python translate_ppt.py input.pptx --api-base https://api.example.com/v1 --api-key your-key --model gpt-4o
        """
    )

    parser.add_argument("input_file", help="Input PPTX file path")
    parser.add_argument("output_file", nargs="?", help="Output PPTX file path (optional)")
    parser.add_argument("--font", default=DEFAULT_FONT, help=f"Target font (default: {DEFAULT_FONT})")
    parser.add_argument("--model", default=DEFAULT_MODEL, help=f"LLM model (default: {DEFAULT_MODEL})")
    parser.add_argument("--api-base", default=DEFAULT_API_BASE, help=f"OpenAI-compatible API base URL (default: {DEFAULT_API_BASE})")
    parser.add_argument("--api-key", default=None,
                       help="API key for the LLM endpoint (optional, not needed for local models)")
    parser.add_argument("--batch-size", type=int, default=DEFAULT_BATCH_SIZE,
                       help=f"Batch size for API calls (default: {DEFAULT_BATCH_SIZE})")
    parser.add_argument("--verbose", "-v", action="store_true", help="Enable verbose logging")

    args = parser.parse_args()

    # Validate input file
    if not os.path.exists(args.input_file):
        print(f"Error: Input file not found: {args.input_file}")
        sys.exit(1)

    if not args.input_file.lower().endswith('.pptx'):
        print("Error: Input file must be a .pptx file")
        sys.exit(1)

    # Determine output path
    output_path = get_output_path(args.input_file, args.output_file)

    # Get API key from args or environment
    api_key = args.api_key or os.environ.get("LLM_API_KEY", "")

    # Soft connectivity check (warning only)
    print(f"Using LLM endpoint: {args.api_base}")
    print(f"Model: {args.model}")
    if not check_api_connectivity(args.api_base):
        print("Warning: Could not verify API connectivity. The endpoint may be unreachable.")
        print("         Will attempt translation anyway...")

    # Load presentation
    print(f"Loading presentation: {args.input_file}")
    try:
        prs = Presentation(args.input_file)
    except Exception as e:
        print(f"Error loading presentation: {e}")
        print("Hint: If the file appears corrupted, try opening and re-saving it in PowerPoint")
        sys.exit(1)

    total_slides = len(prs.slides)
    print(f"Loaded {total_slides} slides")

    # Extract all Chinese text segments
    print("Extracting Chinese text segments...")
    segments = extract_all_texts(prs)

    if not segments:
        print("No Chinese text found in the presentation. Nothing to translate.")
        sys.exit(0)

    print(f"Found {len(segments)} text segments to translate")

    if args.verbose:
        for seg in segments[:10]:  # Show first 10
            print(f"  - Slide {seg.slide_idx + 1}: {seg.text[:50]}...")
        if len(segments) > 10:
            print(f"  ... and {len(segments) - 10} more")

    # Translate segments
    translate_segments(segments, args.model, args.api_base, api_key, args.batch_size, args.verbose)

    # Apply translations
    print("Applying translations to presentation...")
    applied_count = apply_translations(prs, segments, args.font, args.verbose)

    # Font adjustment pass (catch any missed CJK fonts)
    print(f"Adjusting fonts to {args.font}...")
    adjust_all_fonts(prs, args.font)

    # Save output
    print(f"Saving to: {output_path}")
    try:
        prs.save(output_path)
    except Exception as e:
        print(f"Error saving presentation: {e}")
        sys.exit(1)

    # Summary
    print("\n" + "=" * 50)
    print("Translation Complete!")
    print("=" * 50)
    print(f"Slides processed: {total_slides}")
    print(f"Text segments translated: {applied_count}")
    print(f"Output file: {output_path}")
    print("=" * 50)


if __name__ == "__main__":
    main()
