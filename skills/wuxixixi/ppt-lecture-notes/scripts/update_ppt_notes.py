#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT讲师备注自动更新工具 V2
从教案中提取讲师活动，自动写入对应PPT幻灯片备注
"""

import os
import sys
import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime

# 课程目录
DEFAULT_COURSE_DIR = r"G:\own\支教\2026上半年腾讯合作人工智能课"


def install_requirements():
    """安装必要的库"""
    try:
        import pptx
    except ImportError:
        print("正在安装 python-pptx 库...")
        import subprocess
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
        print("安装完成!")


def extract_text_from_docx(docx_path):
    """从docx文件中提取所有文本内容"""
    try:
        with zipfile.ZipFile(docx_path, 'r') as z:
            with z.open('word/document.xml') as f:
                content = f.read()
        
        root = ET.fromstring(content)
        texts = []
        for t in root.iter('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t'):
            if t.text:
                texts.append(t.text)
        return ''.join(texts)
    except Exception as e:
        print(f"读取教案出错: {e}")
        return ""


def find_course_files(course_dir):
    """查找课程文件夹中的教案和课件"""
    courses = []
    
    if not os.path.exists(course_dir):
        print(f"目录不存在: {course_dir}")
        return courses
    
    for item in os.listdir(course_dir):
        item_path = os.path.join(course_dir, item)
        if os.path.isdir(item_path):
            教案 = None
            课件 = None
            
            for f in os.listdir(item_path):
                if '教案' in f and f.endswith('.docx'):
                    教案 = os.path.join(item_path, f)
                if '课件' in f and f.endswith('.pptx'):
                    课件 = os.path.join(item_path, f)
            
            if 教案 and 课件:
                ppt_mtime = os.path.getmtime(课件)
                courses.append({
                    'name': item,
                    '教案': 教案,
                    '课件': 课件,
                    '更新时间': datetime.fromtimestamp(ppt_mtime)
                })
    
    return courses


def get_latest_course(courses):
    """获取最新更新的课程"""
    if not courses:
        return None
    return max(courses, key=lambda x: x['更新时间'])


def update_ppt_notes(教案路径, 课件路径, force=False):
    """更新PPT的讲师备注"""
    from pptx import Presentation

    print(f"\n{'='*50}")
    print(f"处理课程: {os.path.basename(os.path.dirname(课件路径))}")
    print(f"{'='*50}")
    
    # 读取教案内容
    print("读取教案...")
    教案文本 = extract_text_from_docx(教案路径)
    
    if not 教案文本:
        print("无法读取教案内容")
        return None
    
    print("分析教案内容...")
    
    # 打开PPT
    print("打开PPT...")
    try:
        prs = Presentation(课件路径)
    except Exception as e:
        print(f"打开PPT失败: {e}")
        return None
    
    slide_count = len(prs.slides)
    print(f"PPT共有 {slide_count} 张幻灯片")
    
    # 生成输出文件名
    目录 = os.path.dirname(课件路径)
    原名 = os.path.basename(课件路径)
    输出名 = f"课件-更新-{原名}"
    输出路径 = os.path.join(目录, 输出名)
    
    # 从教案中提取讲师活动并生成备注
    讲师活动 = extract_activities_from_教案(教案文本, slide_count)
    
    # 写入备注
    updated_count = 0
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num in 讲师活动:
            # 获取或创建备注
            if not slide.has_notes_slide:
                notes_slide = slide.notes_slide
            else:
                notes_slide = slide.notes_slide
            
            # 检查是否已有备注
            existing_notes = notes_slide.notes_text_frame.text.strip()
            if existing_notes and not force:
                print(f"  Slide {slide_num}: 保留原有备注")
                continue
            
            # 写入新备注
            notes_slide.notes_text_frame.text = 讲师活动[slide_num]
            updated_count += 1
            print(f"  Slide {slide_num}: 已添加备注")
    
    # 保存
    try:
        prs.save(输出路径)
        print(f"\n已保存到: {输出路径}")
        print(f"共更新 {updated_count} 张幻灯片")
    except Exception as e:
        print(f"保存失败: {e}")
        # 尝试保存到桌面
        desktop = os.path.expanduser("~/Desktop")
        输出路径 = os.path.join(desktop, 输出名)
        prs.save(输出路径)
        print(f"已保存到桌面: {输出路径}")
    
    return 输出路径


def extract_activities_from_教案(教案文本, slide_count):
    """从教案文本中提取讲师活动，生成幻灯片备注"""
    notes = {}
    
    # 定义教学阶段对应的幻灯片范围（根据实际教案调整）
    # 这里需要根据具体课程内容进行配置
    
    # 简化版：按段落分割讲师活动
    lines = 教案文本.split('。')
    current_section = ""
    section_start = 1
    
    # 教学阶段定义（从教案中提取）
    stages = []
    stage_keywords = ['一、情景导入', '二、知识讲解', '三、大胆表达', '四、小结', 
                      '五、', '六、', '七、']
    
    for i, line in enumerate(lines):
        for keyword in stage_keywords:
            if keyword in line:
                stages.append((keyword, i))
                break
    
    # 根据阶段生成备注
    if len(stages) >= 2:
        # 按阶段分配幻灯片
        slides_per_stage = slide_count // len(stages)
        for idx, (stage_name, _) in enumerate(stages):
            start_slide = idx * slides_per_stage + 1
            end_slide = (idx + 1) * slides_per_stage if idx < len(stages) - 1 else slide_count + 1
            
            # 提取该阶段的讲师活动
            stage_content = extract_stage_activities(教案文本, stage_name)
            
            # 分配到各幻灯片
            slides_in_stage = end_slide - start_slide
            if slides_in_stage > 0 and stage_content:
                content_per_slide = max(1, len(stage_content) // slides_in_stage)
                for j in range(slides_in_stage):
                    slide_num = start_slide + j
                    if slide_num <= slide_count:
                        content_idx = min(j * content_per_slide, len(stage_content) - 1)
                        notes[slide_num] = stage_name + "\n" + stage_content[content_idx]
    else:
        # 没有明确的阶段划分，使用通用模板
        for i in range(slide_count):
            notes[i + 1] = f"讲师活动：请根据教案内容进行教学。"
    
    return notes


def extract_stage_activities(教案文本, stage_name):
    """提取指定教学阶段的讲师活动"""
    activities = []
    
    # 找到该阶段的位置
    start_idx = 教案文本.find(stage_name)
    if start_idx == -1:
        return activities
    
    # 找到下一个阶段或结束
    next_stages = ['一、', '二、', '三、', '四、', '五、', '六、', '七、']
    next_idx = len(教案文本)
    for ns in next_stages:
        pos = 教案文本.find(ns, start_idx + len(stage_name))
        if pos != -1 and pos < next_idx:
            next_idx = pos
    
    # 提取该阶段内容
    stage_content = 教案文本[start_idx:next_idx]
    
    # 提取讲师活动
    lines = stage_content.split('\n')
    current_activity = []
    in_activity = False
    
    for line in lines:
        if '讲师活动' in line:
            if current_activity:
                activities.append('。'.join(current_activity))
            in_activity = True
            current_activity = []
        elif '学生活动' in line or '设计意图' in line:
            if current_activity:
                activities.append('。'.join(current_activity))
            in_activity = False
        elif in_activity and line.strip():
            current_activity.append(line.strip())
    
    if current_activity:
        activities.append('。'.join(current_activity))
    
    return activities


def main(课程目录=None, force=False):
    """主函数"""
    if not 课程目录:
        课程目录 = DEFAULT_COURSE_DIR
    
    print("=" * 60)
    print("PPT讲师备注自动更新工具 V2")
    print("=" * 60)
    print(f"课程目录: {课程目录}")
    print(f"执行时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    
    # 安装依赖
    install_requirements()
    
    # 查找课程
    print("\n扫描课程文件...")
    courses = find_course_files(课程目录)
    
    if not courses:
        print("未找到任何课程文件（需要包含教案.docx和课件.pptx）")
        return
    
    print(f"\n找到 {len(courses)} 个课程:")
    for i, c in enumerate(courses, 1):
        print(f"  {i}. {c['name']} (PPT更新时间: {c['更新时间'].strftime('%Y-%m-%d %H:%M')})")
    
    # 获取最新更新的课程
    latest = get_latest_course(courses)
    if latest:
        print(f"\n最新更新的课程: {latest['name']}")
        
        # 检查是否已处理过
        目录 = os.path.dirname(latest['课件'])
        原名 = os.path.basename(latest['课件'])
        输出名 = f"课件-更新-{原名}"
        输出路径 = os.path.join(目录, 输出名)
        
        if os.path.exists(输出路径) and not force:
            输出修改时间 = datetime.fromtimestamp(os.path.getmtime(输出路径))
            if 输出修改时间 > latest['更新时间']:
                print(f"\n该课程备注已更新（输出文件更新于 {输出修改时间}）")
                print("使用 --force 参数可强制重新生成")
                return
        
        # 更新备注
        update_ppt_notes(latest['教案'], latest['课件'], force)
    
    print("\n处理完成！")


if __name__ == "__main__":
    课程目录 = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_COURSE_DIR
    force = '--force' in sys.argv
    main(课程目录, force)
