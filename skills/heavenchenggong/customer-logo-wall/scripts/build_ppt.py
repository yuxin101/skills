#!/usr/bin/env python3
"""
build_ppt.py
生成客户 Logo 墙 PPT

用法：
  python3 build_ppt.py \
    --companies companies.json \
    --logos-dir ./logos \
    --output ./output/Logo_Wall.pptx
"""

import argparse
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 配色 ──────────────────────────────────────────────────────
COLORS = {
    "primary":   RGBColor(30, 39, 97),     # #1E2761 深蓝
    "secondary": RGBColor(202, 220, 252),   # #CADCFC 浅蓝
    "accent":    RGBColor(255, 255, 255),   # 白色
    "dark":      RGBColor(15, 22, 49),      # 更深蓝（封面背景）
    "light":     RGBColor(245, 247, 250),   # 浅灰背景
    "text":      RGBColor(30, 39, 97),
    "muted":     RGBColor(107, 114, 128),
}

TIER_COLORS = {
    1: RGBColor(220, 38, 38),    # 红
    2: RGBColor(245, 158, 11),   # 橙
    3: RGBColor(16, 185, 129),   # 绿
    4: RGBColor(107, 114, 128),  # 灰
}

TIER_INFO = {
    1: ("全球顶级知名公司", "Global Top Tier Companies"),
    2: ("行业领先公司",     "Industry Leader Companies"),
    3: ("区域性知名公司",   "Regional Famous Companies"),
    4: ("其他客户",         "Other Customers"),
}


# ── 工具函数 ─────────────────────────────────────────────────
def add_text(slide, text, x, y, w, h, size, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or COLORS["text"]
    p.alignment = align
    return box


def add_rect(slide, x, y, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def find_logo(company, logos_dir):
    """查找公司 logo 文件，返回完整路径或 None"""
    logo_file = company.get('file') or company.get('logo')
    if not logo_file:
        return None
    logo_path = os.path.join(logos_dir, logo_file)
    if os.path.exists(logo_path) and os.path.getsize(logo_path) > 1000:
        return logo_path
    return None


# ── 页面生成 ─────────────────────────────────────────────────
def create_cover(prs, title="Licensed Customers", subtitle="Logo Wall", note=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["dark"]

    add_text(slide, title,    0.5, 1.6, 9, 1.0, 44, True,  COLORS["accent"],    PP_ALIGN.CENTER)
    add_text(slide, subtitle, 0.5, 2.6, 9, 0.8, 36, False, COLORS["secondary"], PP_ALIGN.CENTER)
    if note:
        add_text(slide, note, 0.5, 3.8, 9, 0.5, 16, False, COLORS["muted"],     PP_ALIGN.CENTER)


def create_tier_page(prs, tier_num, page_companies, logos_dir):
    tier_cn, tier_en = TIER_INFO[tier_num]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["light"]

    # 顶部标题栏
    add_rect(slide, 0, 0, 10, 0.8, COLORS["primary"])
    add_text(slide, f"Tier {tier_num} | {tier_cn}", 0.5, 0.15, 9, 0.45, 20, True, COLORS["accent"])
    add_text(slide, tier_en, 0.5, 0.57, 9, 0.22, 10, False, COLORS["secondary"])

    card_w, card_h = 4.3, 1.4
    start_x, start_y = 0.5, 1.0
    gap_x, gap_y = 0.4, 0.2

    for i, c in enumerate(page_companies):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # 卡片底色
        add_rect(slide, x, y, card_w, card_h, COLORS["accent"])
        # 左侧 Tier 色条
        add_rect(slide, x, y, 0.07, card_h, TIER_COLORS[tier_num])

        # Logo
        logo_path = find_logo(c, logos_dir)
        if logo_path:
            try:
                slide.shapes.add_picture(logo_path, Inches(x + 0.2), Inches(y + 0.2), Inches(1.0), Inches(1.0))
            except Exception:
                add_rect(slide, x + 0.2, y + 0.2, 1.0, 1.0, RGBColor(229, 231, 235))
                add_text(slide, "LOGO", x + 0.2, y + 0.55, 1.0, 0.3, 8, True, COLORS["muted"], PP_ALIGN.CENTER)
        else:
            add_rect(slide, x + 0.2, y + 0.2, 1.0, 1.0, RGBColor(229, 231, 235))
            add_text(slide, "LOGO", x + 0.2, y + 0.55, 1.0, 0.3, 8, True, COLORS["muted"], PP_ALIGN.CENTER)

        # 文字
        cn = c.get('cn') or c.get('name', '')
        en = c.get('en', '')
        add_text(slide, cn, x + 1.35, y + 0.25, 2.8, 0.50, 13, True,  COLORS["text"])
        add_text(slide, en, x + 1.35, y + 0.72, 2.8, 0.55,  7, False, COLORS["muted"])


# ── 主函数 ───────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description='生成 Logo 墙 PPT')
    parser.add_argument('--companies', required=True, help='公司 JSON 文件')
    parser.add_argument('--logos-dir', required=True, help='Logo 目录')
    parser.add_argument('--output', required=True, help='输出 PPT 路径')
    parser.add_argument('--title', default='Licensed Customers', help='PPT 标题')
    parser.add_argument('--subtitle', default='Logo Wall', help='PPT 副标题')
    args = parser.parse_args()

    with open(args.companies) as f:
        companies = json.load(f)

    # 自动补 tier（如无 tier 字段则全部归 4）
    for c in companies:
        if 'tier' not in c:
            c['tier'] = 4

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    # 封面
    total = len(companies)
    with_logo = sum(1 for c in companies if find_logo(c, args.logos_dir))
    note = f"共 {total} 家客户  |  按知名度与规模排序"
    create_cover(prs, args.title, args.subtitle, note)

    # 每个 Tier 分页（每页最多 6 个）
    for tier_num in [1, 2, 3, 4]:
        tier_companies = [c for c in companies if c.get('tier') == tier_num]
        for i in range(0, len(tier_companies), 6):
            create_tier_page(prs, tier_num, tier_companies[i:i+6], args.logos_dir)

    os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)
    prs.save(args.output)

    print(f"✅ PPT 已生成: {args.output}")
    print(f"   Logo 覆盖: {with_logo}/{total} 家公司")


if __name__ == '__main__':
    main()
