#!/usr/bin/env python3
"""extract_style.py — Extract precise style data from a PPTX template.

Usage:
    python3 extract_style.py <template.pptx> [--output style_raw.yaml] [--pages 1,5,7,13,26,42]

Outputs a raw style YAML with exact colors, fonts, positions, sizes extracted
from python-pptx. Designed to complement AI visual analysis — the script
provides precision, the AI provides semantic understanding (roles, layout types).
"""

import sys
import os
import argparse
import json
from collections import Counter, defaultdict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False


# ── Helpers ──

def emu_to_inches(emu):
    """Convert EMU to inches, rounded to 2 decimal places."""
    if emu is None:
        return None
    return round(emu / 914400, 2)


def pt_from_emu(emu):
    """Convert EMU font size to points."""
    if emu is None:
        return None
    return round(emu / 12700, 1)


def rgb_to_hex(rgb):
    """Convert RGBColor or theme color to hex string."""
    if rgb is None:
        return None
    if isinstance(rgb, RGBColor):
        return str(rgb).upper()
    return str(rgb).upper()


def extract_color_safe(color_obj):
    """Safely extract color info from a pptx color object."""
    result = {}
    try:
        if color_obj is None:
            return result
        # Try RGB first
        try:
            if color_obj.rgb is not None:
                result["rgb"] = str(color_obj.rgb).upper()
        except (AttributeError, TypeError):
            pass
        # Theme color
        try:
            if color_obj.theme_color is not None:
                result["theme"] = str(color_obj.theme_color)
        except (AttributeError, TypeError):
            pass
        # Brightness
        try:
            if color_obj.brightness is not None and color_obj.brightness != 0:
                result["brightness"] = round(color_obj.brightness, 2)
        except (AttributeError, TypeError):
            pass
    except Exception:
        pass
    return result


def classify_shape_type(shape):
    """Classify shape into our element types."""
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    elif st == MSO_SHAPE_TYPE.MEDIA:
        return "video"
    elif st == MSO_SHAPE_TYPE.TABLE:
        return "table"
    elif st == MSO_SHAPE_TYPE.CHART:
        return "chart"
    elif st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    elif st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text"
    elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        # Check if it has text
        if shape.has_text_frame and shape.text_frame.text.strip():
            return "text_shape"  # shape with text (like tag rectangles)
        return "shape"
    elif st == MSO_SHAPE_TYPE.FREEFORM:
        return "freeform"
    elif st == MSO_SHAPE_TYPE.PLACEHOLDER:
        ph = shape.placeholder_format
        if ph:
            idx = ph.idx
            if idx == 0:
                return "text"  # title placeholder
            elif idx == 1:
                return "text"  # body placeholder
            elif idx in (10, 11, 12, 13, 14, 15, 16, 17, 18):
                return "image"  # picture placeholder
        return "placeholder"
    else:
        return "other"


def extract_text_style(paragraph):
    """Extract text formatting from a paragraph."""
    style = {}
    # Paragraph-level
    if paragraph.alignment is not None:
        style["align"] = str(paragraph.alignment).split("(")[0].lower().replace("pp_align.", "")
    
    # Get dominant run style
    for run in paragraph.runs:
        font = run.font
        if font.size is not None:
            style["fontSize"] = pt_from_emu(font.size)
        if font.name is not None:
            style["fontFace"] = font.name
        if font.bold is not None:
            style["bold"] = font.bold
        if font.italic is not None:
            style["italic"] = font.italic
        if font.underline is not None:
            style["underline"] = bool(font.underline)
        
        color_info = extract_color_safe(font.color)
        if color_info:
            style["color"] = color_info
        break  # Use first run as representative
    
    return style


def extract_fill(shape):
    """Extract fill properties from a shape."""
    fill_info = {}
    try:
        fill = shape.fill
        if fill is None:
            return fill_info
        
        fill_type = fill.type
        if fill_type is not None:
            fill_info["type"] = str(fill_type).replace("MSO_FILL_TYPE.", "").lower()
        
        try:
            if fill.fore_color and fill.fore_color.rgb:
                fill_info["color"] = str(fill.fore_color.rgb).upper()
        except (AttributeError, TypeError):
            pass
        
        try:
            if fill.fore_color and fill.fore_color.theme_color:
                fill_info["theme"] = str(fill.fore_color.theme_color)
        except (AttributeError, TypeError):
            pass
    except Exception:
        pass
    return fill_info


def extract_line_style(shape):
    """Extract line/border properties."""
    line_info = {}
    try:
        ln = shape.line
        if ln is None:
            return line_info
        if ln.width is not None:
            line_info["width"] = round(ln.width / 12700, 1)  # EMU to pt
        color_info = extract_color_safe(ln.color)
        if color_info:
            line_info["color"] = color_info
        if ln.dash_style is not None:
            line_info["dash"] = str(ln.dash_style)
    except Exception:
        pass
    return line_info


def extract_background(slide):
    """Extract slide background properties."""
    bg_info = {}
    try:
        bg = slide.background
        fill = bg.fill
        if fill.type is not None:
            bg_info["type"] = str(fill.type).replace("MSO_FILL_TYPE.", "").lower()
        try:
            if fill.fore_color and fill.fore_color.rgb:
                bg_info["color"] = str(fill.fore_color.rgb).upper()
        except (AttributeError, TypeError):
            pass
        try:
            if fill.fore_color and fill.fore_color.theme_color:
                bg_info["theme"] = str(fill.fore_color.theme_color)
        except (AttributeError, TypeError):
            pass
    except Exception:
        pass
    return bg_info


# ── Main extraction ──

def extract_slide(slide, slide_num):
    """Extract all style data from a single slide."""
    slide_data = {
        "page": slide_num,
        "background": extract_background(slide),
        "elements": [],
    }
    
    for shape in slide.shapes:
        elem = {
            "name": shape.name,
            "type": classify_shape_type(shape),
            "x": emu_to_inches(shape.left),
            "y": emu_to_inches(shape.top),
            "w": emu_to_inches(shape.width),
            "h": emu_to_inches(shape.height),
        }
        
        # Rotation
        if shape.rotation != 0:
            elem["rotation"] = shape.rotation
        
        # Fill (for shapes)
        fill_info = extract_fill(shape)
        if fill_info:
            elem["fill"] = fill_info
        
        # Line/border
        line_info = extract_line_style(shape)
        if line_info:
            elem["line"] = line_info
        
        # Text content & styling
        if shape.has_text_frame:
            tf = shape.text_frame
            text = tf.text.strip()
            if text:
                elem["text_preview"] = text[:80] + ("..." if len(text) > 80 else "")
                elem["paragraph_count"] = len(tf.paragraphs)
                
                # Extract style from each paragraph
                para_styles = []
                for para in tf.paragraphs:
                    ps = extract_text_style(para)
                    if ps:
                        ps["text_preview"] = para.text.strip()[:40]
                        para_styles.append(ps)
                
                if para_styles:
                    elem["text_styles"] = para_styles
                
                # Word wrap
                if tf.word_wrap is not None:
                    elem["word_wrap"] = tf.word_wrap
        
        # Image info
        if elem["type"] == "image":
            try:
                if hasattr(shape, "image"):
                    img = shape.image
                    elem["image_format"] = img.content_type
                    elem["image_size_bytes"] = len(img.blob)
            except Exception:
                pass
        
        # Table info
        if elem["type"] == "table":
            try:
                tbl = shape.table
                elem["rows"] = len(tbl.rows)
                elem["cols"] = len(tbl.columns)
                # Sample first cell style
                if len(tbl.rows) > 0 and len(tbl.columns) > 0:
                    cell = tbl.cell(0, 0)
                    if cell.text.strip():
                        elem["sample_cell"] = cell.text.strip()[:30]
            except Exception:
                pass
        
        slide_data["elements"].append(elem)
    
    return slide_data


def compute_global_stats(slides_data):
    """Aggregate statistics across all analyzed slides."""
    all_colors = Counter()
    all_fonts = Counter()
    all_sizes = Counter()
    all_bg_colors = Counter()
    element_types = Counter()
    
    for sd in slides_data:
        # Background
        bg = sd.get("background", {})
        if "color" in bg:
            all_bg_colors[bg["color"]] += 1
        
        for elem in sd["elements"]:
            element_types[elem["type"]] += 1
            
            # Fill colors
            fill = elem.get("fill", {})
            if "color" in fill:
                all_colors[fill["color"]] += 1
            
            # Text colors, fonts, sizes
            for ps in elem.get("text_styles", []):
                color = ps.get("color", {})
                if "rgb" in color:
                    all_colors[color["rgb"]] += 1
                if "fontFace" in ps:
                    all_fonts[ps["fontFace"]] += 1
                if "fontSize" in ps:
                    all_sizes[ps["fontSize"]] += 1
    
    return {
        "top_colors": dict(all_colors.most_common(15)),
        "top_fonts": dict(all_fonts.most_common(10)),
        "top_font_sizes": dict(all_sizes.most_common(15)),
        "background_colors": dict(all_bg_colors.most_common(5)),
        "element_type_counts": dict(element_types.most_common()),
    }


def main():
    parser = argparse.ArgumentParser(description="Extract style data from PPTX template")
    parser.add_argument("template", help="Path to PPTX template file")
    parser.add_argument("--output", "-o", default=None, help="Output file (default: stdout)")
    parser.add_argument("--pages", "-p", default=None, help="Comma-separated page numbers to analyze (default: all)")
    parser.add_argument("--json", action="store_true", help="Output JSON instead of YAML")
    args = parser.parse_args()
    
    if not os.path.exists(args.template):
        print(f"ERROR: File not found: {args.template}", file=sys.stderr)
        sys.exit(1)
    
    prs = Presentation(args.template)
    
    # Slide size
    slide_w = emu_to_inches(prs.slide_width)
    slide_h = emu_to_inches(prs.slide_height)
    total_slides = len(prs.slides)
    
    # Determine which pages to analyze
    if args.pages:
        pages = [int(p.strip()) for p in args.pages.split(",")]
    else:
        pages = list(range(1, total_slides + 1))
    
    print(f"📐 Slide size: {slide_w}\" × {slide_h}\" ({total_slides} slides)", file=sys.stderr)
    print(f"📊 Analyzing pages: {pages}", file=sys.stderr)
    
    # Extract each slide
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        if i in pages:
            sd = extract_slide(slide, i)
            slides_data.append(sd)
            elem_count = len(sd["elements"])
            text_count = sum(1 for e in sd["elements"] if e["type"] in ("text", "text_shape"))
            img_count = sum(1 for e in sd["elements"] if e["type"] == "image")
            print(f"  Page {i}: {elem_count} elements ({text_count} text, {img_count} images)", file=sys.stderr)
    
    # Global stats
    stats = compute_global_stats(slides_data)
    
    # Build output
    output = {
        "source": os.path.basename(args.template),
        "slide_size": f"{slide_w} x {slide_h} inches",
        "total_slides": total_slides,
        "pages_analyzed": len(slides_data),
        "global_stats": stats,
        "slides": slides_data,
    }
    
    # Serialize
    if args.json or not HAS_YAML:
        text = json.dumps(output, ensure_ascii=False, indent=2, default=str)
    else:
        text = yaml.dump(output, allow_unicode=True, default_flow_style=False, sort_keys=False, width=120)
    
    if args.output:
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"\n✅ Output saved to {args.output}", file=sys.stderr)
    else:
        print(text)


if __name__ == "__main__":
    main()
