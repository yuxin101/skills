#!/usr/bin/env python3
# SECURITY MANIFEST:
# Environment variables accessed: none
# External endpoints called: none
# Local files read: input JSON spec
# Local files written: output PPTX file
"""
Build a PPTX deck from a compact JSON spec.

Usage:
  python scripts/pptx_rebuilder.py spec.json output.pptx

Supports single-slide and multi-slide specs. The goal is not to auto-trace an
image; it is a reusable editable-slide builder that makes it faster to recreate
reference images with native PowerPoint shapes.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any, Dict

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

SHAPE_MAP = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "round_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "chevron": MSO_AUTO_SHAPE_TYPE.CHEVRON,
    "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
    "right_arrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
}

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def rgb(value: str | None, default: str = "000000") -> RGBColor:
    value = (value or default).replace("#", "")
    if len(value) != 6:
        value = default
    return RGBColor.from_string(value.upper())


def inches(v: float) -> int:
    return Inches(v)


def add_textbox(slide, item: Dict[str, Any]):
    tb = slide.shapes.add_textbox(
        inches(item["x"]), inches(item["y"]), inches(item["w"]), inches(item["h"])
    )
    tf = tb.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(item.get("margin", 0.02))
    tf.margin_right = Inches(item.get("margin", 0.02))
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = ALIGN_MAP.get(item.get("align", "center"), PP_ALIGN.CENTER)
    r = p.add_run()
    r.text = item.get("text", "")
    font = r.font
    font.name = item.get("font", "Aptos")
    font.size = Pt(item.get("size", 18))
    font.bold = item.get("bold", False)
    font.color.rgb = rgb(item.get("color"), "000000")
    return tb


def add_shape(slide, item: Dict[str, Any]):
    shape_type = SHAPE_MAP[item["shape"]]
    shp = slide.shapes.add_shape(
        shape_type,
        inches(item["x"]),
        inches(item["y"]),
        inches(item["w"]),
        inches(item["h"]),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(item.get("fill"), "FFFFFF")
    line = item.get("line")
    if line:
        shp.line.color.rgb = rgb(line.get("color"), "000000")
        shp.line.width = Pt(line.get("width", 1))
    else:
        shp.line.fill.background()

    if item.get("text"):
        tf = shp.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = ALIGN_MAP.get(item.get("align", "center"), PP_ALIGN.CENTER)
        r = p.add_run()
        r.text = item["text"]
        font = r.font
        font.name = item.get("font", "Aptos")
        font.size = Pt(item.get("size", 18))
        font.bold = item.get("bold", False)
        font.color.rgb = rgb(item.get("color"), "000000")
    return shp


def apply_background(slide, color: str | None):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(color, "FFFFFF")


def render_slide(prs: Presentation, slide_spec: Dict[str, Any], default_bg: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, slide_spec.get("background", default_bg))
    for item in slide_spec.get("items", []):
        kind = item["kind"]
        if kind == "textbox":
            add_textbox(slide, item)
        elif kind == "shape":
            add_shape(slide, item)
        else:
            raise ValueError(f"Unsupported item kind: {kind}")


def normalize_spec(spec: Dict[str, Any]) -> Dict[str, Any]:
    # Backward compatibility with old single-slide format.
    if "slides" in spec:
        return spec
    presentation = spec.get("presentation") or spec.get("slide") or {}
    return {
        "presentation": {
            "width": presentation.get("width", 13.333),
            "height": presentation.get("height", 7.5),
            "background": presentation.get("background", "FFFFFF"),
        },
        "slides": [
            {
                "background": presentation.get("background", "FFFFFF"),
                "items": spec.get("items", []),
            }
        ],
    }


def main():
    if len(sys.argv) != 3:
        print("Usage: python pptx_rebuilder.py spec.json output.pptx", file=sys.stderr)
        sys.exit(1)

    spec_path = Path(sys.argv[1])
    out_path = Path(sys.argv[2])
    spec = normalize_spec(json.loads(spec_path.read_text()))

    prs = Presentation()
    presentation = spec.get("presentation", {})
    prs.slide_width = Inches(presentation.get("width", 13.333))
    prs.slide_height = Inches(presentation.get("height", 7.5))
    default_bg = presentation.get("background", "FFFFFF")

    # remove default starter slide behavior by creating fresh deck via add_slide only
    if len(prs.slides) > 0:
        pass

    for slide_spec in spec.get("slides", []):
        render_slide(prs, slide_spec, default_bg)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(out_path)


if __name__ == "__main__":
    main()
