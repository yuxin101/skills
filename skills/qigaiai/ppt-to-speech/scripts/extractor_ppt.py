
#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PPT Processor Skill
自动提取 PPT 内容、备注，生成缩略图，并调用 AI 生成演讲稿
"""

import os
import sys
import subprocess
import tempfile
import shutil
import argparse
import platform
import urllib.request
import tarfile
import zipfile
import glob
import stat
import json
import base64
from pathlib import Path

# 第三方库（确保已安装）
try:
    from pptx import Presentation
    from pdf2image import convert_from_path
    from PIL import Image
    import requests
except ImportError as e:
    print(f"缺少 Python 依赖: {e}")
    print("请运行: pip install python-pptx pdf2image Pillow requests")
    sys.exit(1)

# 常量
SCRIPT_DIR = Path(__file__).parent.resolve()
BIN_DIR = SCRIPT_DIR / "bin"
CACHE_DIR = Path.home() / ".cache" / "ppt_skill"
LIBREOFFICE_DIR = CACHE_DIR / "libreoffice"
IMAGE_WIDTH = 800  # 缩略图宽度（像素）

# 平台检测
SYSTEM = platform.system().lower()
ARCH = platform.machine().lower()
IS_WINDOWS = SYSTEM == "windows"
IS_MAC = SYSTEM == "darwin"
IS_LINUX = SYSTEM == "linux"

# LibreOffice 下载地址（阿里云镜像）
LIBREOFFICE_URLS = {
    "linux_x86_64_rpm": "https://mirrors.aliyun.com/libreoffice/stable/25.8.5/rpm/x86_64/LibreOffice_25.8.5_Linux_x86-64_rpm.tar.gz",
    "linux_x86_64_deb": "https://mirrors.aliyun.com/libreoffice/stable/25.8.5/deb/x86_64/LibreOffice_25.8.5_Linux_x86-64_deb.tar.gz",
    "windows_x86_64": "https://mirrors.aliyun.com/libreoffice/stable/25.8.5/win/x86_64/LibreOffice_25.8.5_Win_x86-64.msi",
    "mac_x86_64": "https://mirrors.aliyun.com/libreoffice/stable/25.8.5/mac/x86_64/LibreOffice_25.8.5_MacOS_x86-64.dmg",
}

# Poppler 下载地址（仅 Windows 提供预编译包）
POPPLER_URLS = {
        "windows": "https://github.com/oschwartz10612/poppler-windows/releases/download/v25.12.0-0/Release-25.12.0-0.zip",
}




def install_poppler_linux():
    """Linux 下通过包管理器自动安装 poppler-utils"""
    # 检测包管理器
    if shutil.which("apt-get"):
        pkg_manager = "apt"
        install_cmd = ["apt-get", "install", "-y", "poppler-utils"]
        update_cmd = ["apt-get", "update"]
    elif shutil.which("yum"):
        pkg_manager = "yum"
        install_cmd = ["yum", "install", "-y", "poppler-utils"]
        update_cmd = None
    elif shutil.which("dnf"):
        pkg_manager = "dnf"
        install_cmd = ["dnf", "install", "-y", "poppler-utils"]
        update_cmd = None
    elif shutil.which("pacman"):
        pkg_manager = "pacman"
        install_cmd = ["pacman", "-S", "--noconfirm", "poppler"]
        update_cmd = None
    elif shutil.which("zypper"):
        pkg_manager = "zypper"
        install_cmd = ["zypper", "install", "-y", "poppler-utils"]
        update_cmd = None
    else:
        print("无法识别的包管理器，请手动安装 poppler-utils")
        return False

    # 判断是否需要 sudo
    sudo = []
    if os.geteuid() != 0:
        # 检查 sudo 是否可以无密码运行
        try:
            subprocess.run(["sudo", "-n", "true"], check=True, capture_output=True)
            sudo = ["sudo"]
        except subprocess.CalledProcessError:
            print("需要 sudo 权限安装 poppler-utils，但无法自动获取。")
            return False

    try:
        if update_cmd:
            subprocess.run(sudo + update_cmd, check=True, capture_output=True)
        subprocess.run(sudo + install_cmd, check=True, capture_output=True)
        print(f"已通过 {pkg_manager} 安装 poppler-utils")
        return True
    except subprocess.CalledProcessError as e:
        print(f"通过 {pkg_manager} 安装失败: {e.stderr.decode()}")
        return False


def setup_poppler():
    """确保 poppler 可用，优先使用 BIN_DIR 中的可执行文件，否则尝试自动下载"""
    # 先检查系统中是否有 pdfinfo（poppler 的一部分）
    if shutil.which("pdfinfo"):
        return None  # 系统已有，使用默认 PATH

    # 检查 BIN_DIR 中是否有 pdfinfo
    if IS_WINDOWS:
        pdfinfo_path = BIN_DIR / "pdfinfo.exe"
    else:
        pdfinfo_path = BIN_DIR / "pdfinfo"

    if pdfinfo_path.exists():
        # 确保可执行权限
        if not IS_WINDOWS:
            pdfinfo_path.chmod(pdfinfo_path.stat().st_mode | stat.S_IEXEC)
        return str(BIN_DIR)  # 返回 poppler 目录供 pdf2image 使用

    # 尝试自动安装 poppler
    if IS_WINDOWS:
        print("未检测到 poppler，正在自动下载...")
        try:
            download_poppler_windows()
            return str(BIN_DIR)
        except Exception as e:
            print(f"自动下载 poppler 失败: {e}")
    elif IS_MAC:
        print("未检测到 poppler-utils，请手动安装:")
        print("  brew install poppler")
        print("或下载预编译二进制放到", BIN_DIR)
    elif IS_LINUX:
        print("未检测到 poppler-utils，尝试自动安装...")
        if install_poppler_linux():
            # 安装成功后重新检查
            if shutil.which("pdfinfo"):
                return None
            else:
                print("自动安装后仍未找到 pdfinfo，请手动安装")
        else:
            print("自动安装失败，请手动安装 poppler-utils")
    return None


def download_poppler_windows():
    """下载 Windows 版 poppler 到 BIN_DIR"""
    url = POPPLER_URLS["windows"]
    zip_path = CACHE_DIR / "poppler.zip"
    CACHE_DIR.mkdir(parents=True, exist_ok=True)

    print(f"下载 poppler 从 {url}")
    r = requests.get(url, stream=True)
    r.raise_for_status()
    with open(zip_path, "wb") as f:
        for chunk in r.iter_content(chunk_size=8192):
            f.write(chunk)

    print("解压 poppler...")
    with zipfile.ZipFile(zip_path, "r") as zf:
        zf.extractall(CACHE_DIR)

    # 查找 bin 目录（解压后目录名可能为 poppler-24.08.0/bin 或类似）
    extracted = list(CACHE_DIR.glob("poppler-*/bin"))
    if not extracted:
        # 有时目录结构不同，尝试 Library/bin
        extracted = list(CACHE_DIR.glob("poppler-*/Library/bin"))
    if not extracted:
        raise Exception("未找到 poppler bin 目录")
    poppler_bin = extracted[0]

    # 复制到 BIN_DIR
    BIN_DIR.mkdir(parents=True, exist_ok=True)
    for exe in poppler_bin.glob("*"):
        shutil.copy(exe, BIN_DIR / exe.name)
    print("poppler 已安装到", BIN_DIR)


def setup_libreoffice():
    """确保 LibreOffice 可用，若不可用则自动下载并安装/解压"""
    # 检查系统 PATH 中是否有 soffice
    if shutil.which("soffice"):
        return "soffice"

    # Windows 下检查常见安装目录并临时添加到 PATH
    if IS_WINDOWS:
        lo_candidates = [
            Path("C:/Program Files/LibreOffice/program/soffice.exe"),
            Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
        ]
        for cand in lo_candidates:
            if cand.exists():
                os.environ["PATH"] = str(cand.parent) + os.pathsep + os.environ.get("PATH", "")
                return "soffice"

    # 检查自定义目录
    if IS_WINDOWS:
        lo_exe = LIBREOFFICE_DIR / "program" / "soffice.exe"
    elif IS_MAC:
        lo_exe = LIBREOFFICE_DIR / "Contents" / "MacOS" / "soffice"
    else:
        lo_exe = LIBREOFFICE_DIR / "program" / "soffice"

    if lo_exe.exists():
        os.environ["PATH"] = str(lo_exe.parent) + os.pathsep + os.environ.get("PATH", "")
        return "soffice"

    # 自动下载/安装
    print("未检测到 LibreOffice，正在自动下载...")
    try:
        download_libreoffice()
        return "soffice"
    except Exception as e:
        print(f"自动安装 LibreOffice 失败: {e}")
        print("请手动安装 LibreOffice 并确保 soffice 命令可用。")
        return None


def download_libreoffice():
    """根据平台下载对应 LibreOffice 安装包并解压/提取到用户目录"""
    # 选择下载链接
    if IS_WINDOWS:
        url = LIBREOFFICE_URLS["windows_x86_64"]
        ext = ".msi"
    elif IS_MAC:
        url = LIBREOFFICE_URLS["mac_x86_64"]
        ext = ".dmg"
    elif IS_LINUX:
        if shutil.which("apt"):
            url = LIBREOFFICE_URLS["linux_x86_64_deb"]
        else:
            url = LIBREOFFICE_URLS["linux_x86_64_rpm"]
        ext = ".tar.gz"
    else:
        raise Exception(f"不支持的操作系统: {SYSTEM}")

    print(f"下载 LibreOffice 从 {url}")
    archive_path = CACHE_DIR / f"libreoffice{ext}"
    CACHE_DIR.mkdir(parents=True, exist_ok=True)

    # 下载文件（如果已存在，则跳过）
    if not archive_path.exists():
        r = requests.get(url, stream=True)
        r.raise_for_status()
        with open(archive_path, "wb") as f:
            for chunk in r.iter_content(chunk_size=8192):
                f.write(chunk)
    else:
        print("安装包已存在，跳过下载")

    print("正在准备 LibreOffice...")
    if ext == ".tar.gz":  # Linux
        with tarfile.open(archive_path, "r:gz") as tar:
            tar.extractall(CACHE_DIR)
        # 查找解压后的目录
        extracted = list(CACHE_DIR.glob("LibreOffice_*"))
        if not extracted:
            raise Exception("未找到解压后的 LibreOffice 目录")
        extracted_dir = extracted[0]
        # 移动到 LIBREOFFICE_DIR
        if LIBREOFFICE_DIR.exists():
            shutil.rmtree(LIBREOFFICE_DIR)
        shutil.move(str(extracted_dir), str(LIBREOFFICE_DIR))
        # 将程序目录加入 PATH
        lo_bin = LIBREOFFICE_DIR / "program"
        os.environ["PATH"] = str(lo_bin) + os.pathsep + os.environ.get("PATH", "")
        print("LibreOffice 已解压到", LIBREOFFICE_DIR)
        install_libreoffice_linux(LIBREOFFICE_DIR)

    elif ext == ".msi":  # Windows
        # 如果已经提取过，则直接使用
        lo_exe = LIBREOFFICE_DIR / "program" / "soffice.exe"
        if lo_exe.exists():
            print("LibreOffice 已提取，直接使用")
            os.environ["PATH"] = str(lo_exe.parent) + os.pathsep + os.environ.get("PATH", "")
            return

        # 清空目标目录并重新提取
        if LIBREOFFICE_DIR.exists():
            shutil.rmtree(LIBREOFFICE_DIR)
        LIBREOFFICE_DIR.mkdir(parents=True, exist_ok=True)

        # 使用 msiexec /a 进行管理安装（提取文件，不写入注册表）
        cmd = [
            "msiexec", "/a", str(archive_path),
            "/qn", f"TARGETDIR={str(LIBREOFFICE_DIR)}"
        ]
        print("正在提取 LibreOffice 文件（无需管理员权限）...")
        try:
            subprocess.run(cmd, check=True, capture_output=True, text=True)
        except subprocess.CalledProcessError as e:
            raise Exception(f"提取失败: {e.stderr}")

        # 查找 soffice.exe（可能在 program 子目录下）
        lo_exe = LIBREOFFICE_DIR / "program" / "soffice.exe"
        if not lo_exe.exists():
            # 尝试递归搜索
            possible = list(LIBREOFFICE_DIR.glob("**/soffice.exe"))
            if possible:
                lo_exe = possible[0]
            else:
                raise Exception("提取后未找到 soffice.exe")
        # 将所在目录加入 PATH
        os.environ["PATH"] = str(lo_exe.parent) + os.pathsep + os.environ.get("PATH", "")
        print("LibreOffice 已提取到", LIBREOFFICE_DIR)

    elif ext == ".dmg":  # macOS
        raise Exception("DMG 安装包需要手动挂载安装，请手动安装 LibreOffice 并确保 soffice 命令可用。")


def ppt_to_pdf(ppt_path, output_dir):
    """使用 LibreOffice 将 PPT 转换为 PDF，返回 PDF 文件路径"""
    soffice_cmd = shutil.which("soffice")
    if not soffice_cmd:
        raise Exception("未找到 soffice 命令，无法转换 PDF")

    # 临时目录
    pdf_path = Path(output_dir) / "presentation.pdf"
    cmd = [
        soffice_cmd,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(ppt_path)
    ]
    print("正在转换 PPT 到 PDF...")
    subprocess.run(cmd, check=True, capture_output=True, text=True)

    # LibreOffice 生成的文件名与输入相同但后缀 .pdf
    generated_pdf = Path(output_dir) / (Path(ppt_path).stem + ".pdf")
    if generated_pdf.exists():
        shutil.move(str(generated_pdf), str(pdf_path))
    else:
        # 尝试查找其他可能名称
        for f in Path(output_dir).glob("*.pdf"):
            shutil.move(str(f), str(pdf_path))
            break

    if not pdf_path.exists():
        raise Exception("PDF 转换失败")
    return pdf_path


def extract_ppt_content(ppt_path):
    """使用 python-pptx 提取每页内容"""
    prs = Presentation(ppt_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        # 提取标题
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
        else:
            # 尝试从第一个文本框获取标题
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    title = shape.text.strip()
                    break

        # 提取正文（除标题外的文本）
        body_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    body_texts.append(text)

        # 提取备注
        notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes = notes_frame.text.strip()

        slides_data.append({
            "page": idx,
            "title": title,
            "body": "\n\n".join(body_texts),
            "notes": notes
        })

    return slides_data


def generate_thumbnails(pdf_path, thumb_dir, poppler_path=None):
    """从 PDF 生成每页缩略图"""
    print("正在生成缩略图...")
    images = convert_from_path(
        pdf_path,
        size=(IMAGE_WIDTH, None),
        poppler_path=poppler_path
    )
    thumb_dir = Path(thumb_dir)
    thumb_dir.mkdir(parents=True, exist_ok=True)

    for i, img in enumerate(images, start=1):
        img_path = thumb_dir / f"page_{i:03d}.png"
        img.save(img_path, "PNG")
    return len(images)


def install_libreoffice_linux(LIBREOFFICE_DIR):
    """安装 LibreOffice RPM 包，并确保 soffice 命令在当前进程中可用"""
    debs_dir = os.path.join(LIBREOFFICE_DIR, "DEBS")
    rpms_dir = os.path.join(LIBREOFFICE_DIR, "RPMS")

    if os.path.exists(debs_dir) and os.listdir(debs_dir):
        # 处理 DEB 包 (Ubuntu/Debian 系统)
        return install_libreoffice_deb(debs_dir)
    elif os.path.exists(rpms_dir) and os.listdir(rpms_dir):
        # 处理 RPM 包 (CentOS/RHEL/Fedora 系统)
        return install_libreoffice_rpm(rpms_dir)
    else:
        print(f"未找到安装包目录: 检查 {debs_dir} 或 {rpms_dir}")
        return False


def install_libreoffice_deb(debs_dir):
    """安装 DEB 格式的 LibreOffice"""
    deb_files = glob.glob(os.path.join(debs_dir, "*.deb"))
    if not deb_files:
        print(f"未找到 DEB 包: {debs_dir}")
        return False

    print(f"找到 {len(deb_files)} 个 DEB 包，开始安装...")

    # 构建安装命令
    cmd = ["dpkg", "-i"] + deb_files

    # 非 root 用户添加 sudo
    if os.geteuid() != 0:
        cmd = ["sudo"] + cmd

    try:
        subprocess.run(cmd, check=True)
        print("DEB 包安装完成")
    except subprocess.CalledProcessError as e:
        print(f"安装失败: {e}")
        return False

    # 查找 soffice 安装位置
    # DEB 包通常安装到 /opt/libreoffice*/program/
    return find_and_setup_soffice("/opt")


def install_libreoffice_rpm(rpms_dir):
    """安装 RPM 格式的 LibreOffice (原函数逻辑)"""
    rpm_files = glob.glob(os.path.join(rpms_dir, "*.rpm"))
    if not rpm_files:
        print(f"未找到 RPM 包: {rpms_dir}")
        return False

    # 选择包管理器
    if shutil.which("dnf"):
        cmd = ["dnf", "install", "-y"] + rpm_files
    elif shutil.which("yum"):
        cmd = ["yum", "localinstall", "-y"] + rpm_files
    else:
        print("未检测到 dnf 或 yum")
        return False

    # 非 root 用户添加 sudo
    if os.geteuid() != 0:
        cmd = ["sudo"] + cmd

    try:
        subprocess.run(cmd, check=True)
        print("RPM 包安装完成")
    except subprocess.CalledProcessError as e:
        print(f"安装失败: {e}")
        return False

    return find_and_setup_soffice("/opt")


def find_and_setup_soffice(search_base="/opt"):
    """
    查找 soffice 可执行文件并设置 PATH
    """
    # 查找可能的 soffice 位置
    possible_paths = [
        "/usr/bin/soffice",
        "/usr/local/bin/soffice",
    ]

    # 搜索 /opt/libreoffice*/program/soffice
    opt_paths = glob.glob(f"{search_base}/libreoffice*/program/soffice")
    possible_paths.extend(opt_paths)

    for soffice_path in possible_paths:
        if os.path.isfile(soffice_path) and os.access(soffice_path, os.X_OK):
            # 创建软链接到 /usr/local/bin (如果还没有)
            if not os.path.exists("/usr/local/bin/soffice"):
                try:
                    os.symlink(soffice_path, "/usr/local/bin/soffice")
                    print(f"已创建软链接: {soffice_path} -> /usr/local/bin/soffice")
                except OSError:
                    pass

            # 添加到当前进程 PATH
            bin_dir = os.path.dirname(soffice_path)
            os.environ["PATH"] = bin_dir + os.pathsep + os.environ.get("PATH", "")
            print(f"soffice 已找到并配置: {soffice_path}")
            return True

    print("警告：未找到 soffice 可执行文件")
    return False



def main():
    parser = argparse.ArgumentParser(description="PPT 转演讲稿 Skill")
    parser.add_argument("input", help="输入 PPT 文件路径")
    parser.add_argument("--output", help="输出 Markdown 文件路径（默认与输入同目录，后缀 .md）")
    args = parser.parse_args()

    input_path = Path(args.input).resolve()
    if not input_path.exists():
        print(f"错误: 文件不存在 {input_path}")
        sys.exit(1)

    if input_path.suffix.lower() != ".pptx":
        print("错误: 仅支持 .pptx 格式")
        sys.exit(1)

    # 设置输出路径
    if args.output:
        output_path = Path(args.output).resolve()
    else:
        output_path = input_path.parent / (input_path.stem + ".md")

    # 创建临时工作目录
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir)

        # 1. 安装/检查依赖
        poppler_path = setup_poppler()
        if not setup_libreoffice():
            sys.exit(1)

        # 2. 提取 PPT 内容
        slides_data = extract_ppt_content(input_path)

        # 3. 转换为 PDF 并生成缩略图
        try:
            pdf_file = ppt_to_pdf(input_path, tmp_path)
            thumb_dir = input_path.parent / (input_path.stem + "_thumbnails")
            num_pages = generate_thumbnails(pdf_file, thumb_dir, poppler_path)
            print(f"生成 {num_pages} 张缩略图")
        except Exception as e:
            print(f"生成缩略图失败: {e}")
            thumb_dir = None

    print("处理完成！")


if __name__ == "__main__":
    main()
